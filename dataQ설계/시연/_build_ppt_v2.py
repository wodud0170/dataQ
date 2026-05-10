"""
Narae DataQ 소개서 PPT 생성기 v3 (2026-05-10).

컨셉: 고객 요구사항(Need) ↔ DataQ 답변(Answer) 두 컬럼 매핑.
화면별 상세는 라이브 시연으로 진행 — PPT 는 가치 제안에 집중.

디자인:
  - Cover / Section / Closing: 다크 인디고 + 시안 액센트
  - Content: 화이트 + 좌(Need 회색 카드) ↔ 우(Answer 화이트 카드 + violet 헤더)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "Narae_DataQ_소개서_2026-05-10.pptx")

# ========== 컬러 ==========
DARK_BG    = RGBColor(0x0E, 0x0E, 0x2C)
WHITE_BG   = RGBColor(0xFA, 0xFA, 0xFA)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK  = RGBColor(0x0F, 0x17, 0x2A)
TEXT_BODY  = RGBColor(0x33, 0x41, 0x55)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
CARD_LINE  = RGBColor(0xE2, 0xE8, 0xF0)
NEED_BG    = RGBColor(0xF1, 0xF5, 0xF9)
NEED_LBL   = RGBColor(0x47, 0x55, 0x69)
VIOLET     = RGBColor(0x7C, 0x3A, 0xED)
CYAN       = RGBColor(0x06, 0xB6, 0xD4)
INDIGO     = RGBColor(0x43, 0x38, 0xCA)
ANSWER_HD  = RGBColor(0x4F, 0x46, 0xE5)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_slide(): return prs.slides.add_slide(BLANK)
def fill(s, c): s.fill.solid(); s.fill.fore_color.rgb = c
def line_off(s): s.line.fill.background()


def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    fill(s, color); line_off(s); return s


def rrect(slide, l, t, w, h, color, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    fill(s, color)
    if line_color is None:
        line_off(s)
    else:
        s.line.color.rgb = line_color
        s.line.width = Pt(0.5)
    return s


def text(slide, l, t, w, h, txt, *, size=18, bold=False, color=None,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Pretendard'):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = font
    if color is not None: r.font.color.rgb = color
    return box


def bullets(slide, l, t, w, h, items, *, size=14, color=TEXT_BODY,
            line_space=1.45, font='Pretendard', bullet='·'):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.line_spacing = line_space
        r = p.add_run(); r.text = f"{bullet}  {item}"
        r.font.size = Pt(size); r.font.name = font
        r.font.color.rgb = color


def page_no(slide, n, total):
    text(slide, Inches(12.6), Inches(7.15), Inches(0.7), Inches(0.3),
         f"{n} / {total}", size=9, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


# ========== 슬라이드 빌더 ==========
def cover():
    s = add_slide()
    rect(s, 0, 0, SW, SH, DARK_BG)
    blob = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(-2.0),
                              Inches(7.5), Inches(7.5))
    fill(blob, INDIGO); line_off(blob); blob.fill.transparency = 0.5
    blob2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2.5), Inches(4.5),
                               Inches(5.5), Inches(5.5))
    fill(blob2, VIOLET); line_off(blob2); blob2.fill.transparency = 0.6
    chip = rrect(s, Inches(0.8), Inches(0.9), Inches(3.2), Inches(0.45), VIOLET)
    text(s, Inches(0.8), Inches(0.9), Inches(3.2), Inches(0.45),
         "DATA QUALITY PLATFORM", size=12, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(0.8), Inches(2.0), Inches(11), Inches(1.2),
         "Narae DataQ", size=64, bold=True, color=WHITE)
    text(s, Inches(0.8), Inches(3.1), Inches(11), Inches(0.8),
         "데이터 표준화·품질 관리 통합 솔루션", size=28, color=CYAN)
    text(s, Inches(0.8), Inches(4.1), Inches(11), Inches(0.6),
         "고객의 요구를, 기능 한 줄로 답합니다.", size=18, color=TEXT_MUTED)
    text(s, Inches(0.8), Inches(6.6), Inches(8), Inches(0.4),
         "2026.05.10  |  소개서",
         size=12, color=TEXT_MUTED)


def overview(total):
    s = add_slide()
    rect(s, 0, 0, SW, SH, WHITE_BG)
    rect(s, 0, 0, Inches(0.12), Inches(3.75), VIOLET)
    rect(s, 0, Inches(3.75), Inches(0.12), Inches(3.75), CYAN)
    text(s, Inches(0.6), Inches(0.55), Inches(12), Inches(0.4),
         "Narae DataQ — 한 장 소개", size=11, color=TEXT_MUTED)
    text(s, Inches(0.6), Inches(0.95), Inches(12), Inches(0.8),
         "데이터 라이프사이클 전체를, 한 솔루션에서.", size=28, bold=True, color=TEXT_DARK)
    rect(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.02), CARD_LINE)

    # 5 단계 가로 카드
    steps = [
        ("표준 사전",  "단어 · 용어 · 도메인",       VIOLET),
        ("모델 수집",  "DBMS · XMI 자동 임포트",     INDIGO),
        ("표준 진단",  "컬럼 표준 적합성 검사",      ANSWER_HD),
        ("품질 진단",  "값 프로파일 + 룰 검증",      CYAN),
        ("자동화",     "스케줄 + 회귀 + 자동 표준화", RGBColor(0x10, 0xB9, 0x81)),
    ]
    card_w = Inches(2.34)
    gap = Inches(0.07)
    start_x = Inches(0.6)
    for i, (k, v, c) in enumerate(steps):
        x = start_x + (card_w + gap) * i
        card = rrect(s, x, Inches(2.3), card_w, Inches(2.0), WHITE, line_color=CARD_LINE)
        rect(s, x, Inches(2.3), card_w, Inches(0.16), c)
        text(s, x + Inches(0.2), Inches(2.55), card_w - Inches(0.4), Inches(0.5),
             k, size=18, bold=True, color=TEXT_DARK)
        text(s, x + Inches(0.2), Inches(3.15), card_w - Inches(0.4), Inches(1.0),
             v, size=12, color=TEXT_BODY)

    # 키 메시지 한 줄
    box = rrect(s, Inches(0.6), Inches(4.7), Inches(12.0), Inches(2.1), WHITE,
                line_color=CARD_LINE)
    text(s, Inches(0.9), Inches(4.95), Inches(11.5), Inches(0.5),
         "본 발표는 화면 위주 설명이 아닙니다.", size=18, bold=True, color=ANSWER_HD)
    bullets(s, Inches(0.9), Inches(5.5), Inches(11.5), Inches(1.3), [
        "RFP SFR-22 의 A · C 두 영역(메타 수집·관리 / 메타시스템 연계)이 DataQ 의 책임 영역입니다.",
        "B 영역(DB 모델링 관리)은 ERwin 도입으로 충족되며, DataQ 는 그 결과물을 받아 표준화·진단으로 연결합니다.",
        "여기에 자동 표준화 · 데이터 품질 진단 · 거버넌스 · 자동화 회귀라는 4개 차별화 가치를 더합니다.",
        "구체 화면·동작은 본 PPT 직후 라이브 시연으로 직접 확인해주십시오.",
    ], size=13)
    page_no(s, 2, total)


def section(idx, title, subtitle, total, page_n):
    s = add_slide()
    rect(s, 0, 0, SW, SH, DARK_BG)
    blob = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5),
                              Inches(6), Inches(6))
    fill(blob, VIOLET); line_off(blob); blob.fill.transparency = 0.7
    text(s, Inches(0.8), Inches(0.9), Inches(3), Inches(0.5),
         f"CHAPTER {idx:02d}", size=14, bold=True, color=CYAN)
    text(s, Inches(0.8), Inches(2.5), Inches(11), Inches(1.5),
         title, size=48, bold=True, color=WHITE)
    text(s, Inches(0.8), Inches(4.2), Inches(11), Inches(0.6),
         subtitle, size=18, color=TEXT_MUTED)
    page_no(s, page_n, total)


def need_answer(n, total, header, pairs):
    """
    pairs: [(need_text, [answer_lines]), ...]  — 카드 1~3쌍 권장
    """
    s = add_slide()
    rect(s, 0, 0, SW, SH, WHITE_BG)
    rect(s, 0, 0, Inches(0.12), Inches(3.75), VIOLET)
    rect(s, 0, Inches(3.75), Inches(0.12), Inches(3.75), CYAN)
    text(s, Inches(0.6), Inches(0.55), Inches(12), Inches(0.4),
         "Narae DataQ — 소개서", size=11, color=TEXT_MUTED)
    text(s, Inches(0.6), Inches(0.95), Inches(12), Inches(0.8),
         header, size=26, bold=True, color=TEXT_DARK)
    rect(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.02), CARD_LINE)

    # 카드 행 — 최대 3 쌍
    n_pairs = len(pairs)
    top_y = Inches(2.15)
    avail_h = Inches(5.0)  # 6.85 - 2.15
    gap_y = Inches(0.15)
    each_h_emu = (avail_h.emu - gap_y.emu * (n_pairs - 1)) // n_pairs
    each_h = each_h_emu

    for i, (need, answer_lines) in enumerate(pairs):
        y = top_y + (each_h + gap_y) * i
        # NEED (좌측, 회색)
        rrect(s, Inches(0.6), y, Inches(4.0), each_h, NEED_BG)
        text(s, Inches(0.8), y + Inches(0.15), Inches(3.6), Inches(0.35),
             "고객 요구", size=10, bold=True, color=NEED_LBL)
        text(s, Inches(0.8), y + Inches(0.5), Inches(3.6), each_h - Inches(0.6),
             need, size=14, color=TEXT_DARK)
        # ANSWER (우측, 화이트 + violet 헤더 라벨)
        rrect(s, Inches(4.75), y, Inches(7.85), each_h, WHITE, line_color=CARD_LINE)
        text(s, Inches(4.95), y + Inches(0.15), Inches(7.5), Inches(0.35),
             "DataQ 답변", size=10, bold=True, color=ANSWER_HD)
        bullets(s, Inches(4.95), y + Inches(0.5), Inches(7.5), each_h - Inches(0.6),
                answer_lines, size=12, line_space=1.4)
    page_no(s, n, total)


def closing(total):
    s = add_slide()
    rect(s, 0, 0, SW, SH, DARK_BG)
    blob = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-3), Inches(-2),
                              Inches(8), Inches(8))
    fill(blob, INDIGO); line_off(blob); blob.fill.transparency = 0.55
    blob2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(3),
                               Inches(7), Inches(7))
    fill(blob2, CYAN); line_off(blob2); blob2.fill.transparency = 0.7
    text(s, Inches(0.8), Inches(2.4), Inches(11), Inches(1.0),
         "지금부터 라이브 시연을 시작합니다.", size=36, bold=True, color=WHITE)
    text(s, Inches(0.8), Inches(3.4), Inches(11), Inches(0.6),
         "각 카테고리의 답변은 실제 화면으로 직접 확인 가능합니다.", size=18, color=CYAN)
    text(s, Inches(0.8), Inches(6.6), Inches(11), Inches(0.4),
         "Narae DataQ  ·  Q&A 환영합니다",
         size=12, color=TEXT_MUTED)
    page_no(s, total, total)


# ========== 콘텐츠 ==========
CHAPTERS = [
    # (idx, title, subtitle, [page]: list of (header, [(need, answer_lines), ...]))

    # === Chapter 1 — RFP 요구사항 A: 메타데이터 수집·관리 SW ===
    ("01", "RFP 요구사항 A — 메타데이터 수집·관리",
     "스키마·테이블·컬럼·인덱스·제약조건 자동 수집 + 변경 이력",
     [
         ("DB 구조 정보 자동 수집", [
             ("발주처가 운영하는 다양한 DBMS 의 메타데이터를 한 시스템에서 통합 수집·관리하고 싶습니다.",
              [
                  "Oracle (SID / Service Name) · PostgreSQL · Cubrid · MariaDB / MySQL 다중 DBMS 지원",
                  "데이터소스 등록 → 스키마 다중 선택 → 카탈로그 직접 접속 자동 수집",
                  "스키마 / 테이블 / 컬럼 / 인덱스 / 제약조건 — 5종 객체 일괄 수집",
                  "수집 진행 상황 WebSocket 실시간 push, 결과 카운트 표시",
              ]),
             ("ERwin 같은 모델링 도구 결과물도 동일 시스템에 적재되어야 합니다.",
              [
                  "OMG **XMI 2.1** 표준 임포트 — ERwin export 파일 그대로 업로드",
                  "한글 alias(ownedComment) 자동 인식 — 한글명 별도 입력 불필요",
                  "테이블 · 컬럼 · 인덱스 · 제약조건 · FK 한 번에 적재",
              ]),
         ]),
         ("변경 이력 자동 추적", [
             ("운영 DB 구조가 변경되었는지 자동으로 감지하고 이력으로 남기고 싶습니다.",
              [
                  "구조 변경 진단 — 수집 스냅샷 vs 현재 DB 비교, ADDED/MODIFIED/DELETED 자동 감지",
                  "TB_STRUCT_DIAG_HISTORY (이력) + DETAIL/INDEX_DETAIL/CONSTRAINT_DETAIL (변경 상세)",
                  "결과 화면 3탭 — 컬럼 / 인덱스 / 제약조건. 탭별 엑셀 다운로드",
                  "스케줄 등록 시 매일 새벽 자동 진단 가능 — 변경 보고서만 받으면 됨",
              ]),
             ("표준 사전·모델 변경에 대한 사용자 단위 감사 추적이 필요합니다.",
              [
                  "TB_CHANGE_HISTORY — 단어/용어/도메인/코드 등록·수정·삭제 모든 행위 기록",
                  "사용자 / 시각 / before·after 값 영구 보관 (모델 삭제 후에도 유지)",
                  "변경 이력 화면에서 검색·필터·상세 조회",
              ]),
         ]),
     ]),

    # === Chapter 2 — RFP 요구사항 B: ERwin 영역 (DataQ 직접 충족 영역 아님) ===
    ("02", "RFP 요구사항 B — 모델링 관리 (ERwin 영역)",
     "DataQ 는 ERwin 모델을 임포트해 표준화·DDL 라인으로 연계",
     [
         ("RFP B 는 모델링 도구 영역", [
             ("ERD 편집·시각화·대용량 처리는 DataQ 가 직접 담당하나요?",
              [
                  "**RFP B 의 \"DB 모델링 관리 SW\" 는 ERwin · DA# 등 전문 모델링 도구 영역**",
                  "발주처는 ERwin 라이선스 도입 + 환경 설치 지원으로 RFP B 를 충족",
                  "DataQ 는 ERwin 의 결과물을 받아 표준 검증·DDL·메타 관리로 연결하는 역할",
                  "(SMETA 등 경쟁 메타 솔루션도 동일 — ERD 자체 뷰어는 보유하지 않음)",
              ]),
             ("DataQ 는 ERwin 과 어떻게 연계되나요?",
              [
                  "OMG **XMI 2.1** 표준 임포트 — ERwin export 파일 그대로 업로드",
                  "한글 alias(ownedComment) 자동 인식",
                  "임포트된 모델은 표준 진단 / 자동 표준화 / DDL 다운로드 흐름으로 즉시 이어짐",
                  "(자세히는 RFP C — 메타시스템 연계)",
              ]),
         ]),
     ]),

    # === Chapter 3 — RFP 요구사항 C: 메타시스템 연계 모듈 ===
    ("03", "RFP 요구사항 C — 메타시스템 연계",
     "표준 포맷 추출 · 변환 · 적재 (ETL 패턴)",
     [
         ("표준 포맷 추출 / 적재", [
             ("외부 메타시스템·모델링 도구와 데이터 모델을 주고받고 싶습니다.",
              [
                  "**추출(Export)** — 테이블/컬럼/인덱스/제약조건/단어/용어/도메인 전체 Excel 다운로드 + DDL SQL 생성",
                  "**적재(Load/Import)** — 단어/용어/도메인/코드/모델 Excel 일괄 업로드 + ERwin XMI 임포트",
                  "RFP 한국어 파일명 RFC 5987 (UTF-8 + ASCII fallback) — 다운로드 호환성 확보",
                  "REST API 전체 노출 — `/api/dm/...` 외부 시스템 연동 가능",
              ]),
             ("연계 작업 결과를 추적하고 싶습니다.",
              [
                  "Excel 업로드 미리보기 → 커밋 패턴 — 잘못 올린 파일은 미리보기 단계에서 차단",
                  "임포트 결과 카운트 (성공/실패/경고) 즉시 표시",
                  "TB_CHANGE_HISTORY 에 자동 기록 — 누가 언제 어떤 파일로 적재했는지 추적",
              ]),
         ]),
     ]),

    # === Chapter 4 — DataQ 차별화: 표준 사전 통합 + 자동 표준화 ===
    ("04", "DataQ 차별화 ① — 표준 사전 + 자동 표준화",
     "한국 표준 직접 반영 + DP 점수 단어 분리 + 양방향 표준화",
     [
         ("표준 사전 통합", [
             ("표준 단어·용어·도메인이 엑셀과 사내 위키에 흩어져 있어 어디가 최신인지 모르겠습니다.",
              [
                  "단어 · 용어 · 도메인 · 코드 · 도메인 분류 — 5종을 한 시스템에 통합",
                  "**행안부 공통 표준 단어 약 3,300건 기본 탑재** — 사내 단어를 그 위에 누적",
                  "한국 표준 직접 반영 — 형식단어 / 분류어 / 한글 도메인 분류 트리",
                  "변경 이력 자동 누적 — 누가 언제 무엇을 (before/after) 변경했는지",
              ]),
         ]),
         ("자동 표준화 — 수동 매핑 최소화", [
             ("용어 만들 때마다 단어 분리·영문약어 매핑이 너무 손이 갑니다.",
              [
                  "한글 용어 입력 → 1초 디바운스 후 자동 분석",
                  "**DP 점수 + 동의어 cascade + 미신뢰 토큰 재검색** — 정확한 단어 분리",
                  "영문약어 자동 합성 (예: 회원전화번호 → MBR_TEL_NO)",
                  "마지막 토큰이 형식단어(`CD` 등) → 도메인 유형 자동 cascade + 코드 picker",
                  "미등록 단어 — 같은 모달 안에서 인라인 등록 (모달 안 모달 없음)",
              ]),
             ("이미 운영 중인 한글 컬럼명 수백 개를 표준화하려니 일일이 매핑이 어렵습니다.",
              [
                  "한글 컬럼명 textarea 다중 입력 → 일괄 분석 → REGISTERED / AUTO / PARTIAL / FAILED 4종 분기",
                  "컬럼별 [한글명 기준 표준화] / [영문명 기준 표준화] **양방향**",
                  "**dryRun 정책** — 표준화 결과는 그리드 셀에만 채우고 DB는 그대로. [저장] 시점에 일괄 반영",
                  "잘못 클릭해도 새로고침으로 즉시 복귀 — 운영 안전",
              ]),
         ]),
     ]),

    # === Chapter 5 — DataQ 차별화: 데이터 품질 진단 ===
    ("05", "DataQ 차별화 ② — 데이터 품질 진단",
     "도메인 룰 + 값 프로파일 + 업무 규칙 + 시계열 통계 통합",
     [
         ("표준만이 아닌 \"값\" 검증", [
             ("타입·길이 표준에 맞아도 실제 값에 NULL · 범위 초과 · 형식 오류가 섞여 있습니다.",
              [
                  "도메인 룰 8종 — NOT_NULL / RANGE / LENGTH / REGEX / ENUM / UNIQUE / REFERENCE / COMPARE",
                  "각 유형 전용 입력 위젯 — RANGE min/max + 정수 토글, REGEX 실시간 매칭 테스트, ENUM 칩 입력",
                  "분류 단위 정의 → 같은 분류 도메인 모두에 자동 적용",
                  "위반 샘플 100건 자동 추출 — PK + 위반값 즉시 확인",
              ]),
             ("이메일·주민번호 같은 표준 정규식을 룰마다 직접 짜기 번거롭습니다.",
              [
                  "**룰 카탈로그** — 시스템 기본 50+ 종 (이메일 · 전화 · 주민번호 등) 사전 등록",
                  "한 번 클릭으로 매핑 (Import) 또는 [복사] 후 사용자 정의로 fork 편집",
                  "카탈로그 검색 + 파라미터 / 설명 한눈에",
              ]),
         ]),
         ("프로파일링 + 시계열", [
             ("값 분포 (NULL 비율, 고유값 수 등) 통계와 시간에 따른 추이도 같이 보고 싶습니다.",
              [
                  "값 프로파일링 — 컬럼 단위 분포·이상값 진단",
                  "샘플링 1만건 / 10% / 100% 선택 + 도메인 분류·테이블·컬럼 다중 필터",
                  "**진단 통계** — 모델/컬럼 단위 적합률 시계열 (최근 30회) ApexChart 라인 차트",
                  "6개 품질 화면 모두 .catch + .finally 보강 — 무한 로딩 차단",
              ]),
         ]),
     ]),

    # === Chapter 6 — DataQ 차별화: 거버넌스 + 운영 안전 ===
    ("06", "DataQ 차별화 ③ — 거버넌스 + 운영 안전 장치",
     "권한 / 승인 / cascade / 인라인 편집 / dryRun",
     [
         ("권한 + 승인", [
             ("아무나 표준 단어를 등록할 수 있으면 사전이 오염됩니다.",
              [
                  "관리자 / 일반 사용자 권한 2단계 — 관리자 그룹 메뉴는 일반 사용자에게 미노출",
                  "일반 사용자 등록 → APRV_YN='N' → 관리자 [승인]/[반려] 인라인 처리",
                  "**단어 선승인** 정책 — 미승인 단어 포함 용어는 승인 거부",
                  "**cascade 반려** — 단어 반려 시 연관 미승인 용어 동시 반려",
              ]),
         ]),
         ("그리드 운영 편의 + 안전 장치", [
             ("ERwin 으로 만든 모델을 DataQ 에서 운영하다 보면 손이 많이 가는 작업이 생깁니다.",
              [
                  "그리드 모든 셀 **인라인 편집** — 영문명 / 타입(드롭다운+자유입력) / 길이 / 소수점 / 순서 / 한글명 / NULL·PK·FK / 디폴트",
                  "엑셀에서 한 컬럼 복사 → Ctrl+V 로 같은 필드의 다음 행에 자동 분배 (멀티 paste)",
                  "추가 대상 테이블 변경 시 검색 필터(소유자/테이블) 자동 세팅",
                  "영문명 자동 대문자, 컬럼 순서 자동 1·2·3 재배열, 영문명/한글명 중복 사전 차단",
              ]),
             ("컬럼 영문명 변경, 테이블명 변경, 모델 삭제 시 사고가 자주 납니다.",
              [
                  "**컬럼 영문명 변경 → INDEX / CONSTRAINT / FK_PARENT_ATTR cascade rename** 자동",
                  "**테이블 물리명 변경 → OBJ_NM + ATTR + INDEX + CONSTRAINT + REF_TABLE 5단계 cascade** + 영향 카운트 미리보기",
                  "**모델 삭제 → 연관 진단 스케줄 자동 비활성** + 경고 confirm + 비활성 건수 안내",
                  "**진단 제외 관리** — 임시·폐기 테이블 ON/OFF + 사유 기록 (표준/구조/품질 3종)",
                  "**dryRun 표준화** — DB 즉시 반영 안 하고 그리드만 갱신, [저장] 눌러야 적용",
              ]),
         ]),
     ]),

    # === Chapter 7 — DataQ 차별화: 자동화 ===
    ("07", "DataQ 차별화 ④ — 자동화 + 신뢰도",
     "스케줄 + 자동 회귀 약 100여 케이스",
     [
         ("진단 스케줄 자동 실행", [
             ("매주 월요일 아침 진단을 사람이 직접 클릭하고 있습니다.",
              [
                  "SIMPLE 스케줄 (일/주/월 + 시각) 또는 CRON 표현식 자유 등록",
                  "진단 유형 — STANDARD / STRUCT / BOTH 동시 가능",
                  "[즉시 실행] 버튼으로 수동 트리거",
                  "**동시 실행 방어** — 같은 모델·유형 진행 중이면 자동 SKIP + 로그",
                  "실행 이력 자동 누적 — 시간/모델/유형/상태/소요시간/결과 카운트",
              ]),
         ]),
         ("자동화 회귀 — 신뢰도 보증", [
             ("새 기능을 추가할 때마다 기존 화면이 깨지는 일이 잦습니다.",
              [
                  "Selenium 통합 테스트 약 100여 케이스 — 한 명령어로 일괄 실행",
                  "데이터 표준 사전 / 데이터 모델 / 표준·구조 진단 / 데이터 품질 6개 화면 / 스케줄 / 용어 등록 v2 30 케이스",
                  "테스트마다 cleanup hook — DB 폴루션 자동 정리, 좀비 브라우저 종료",
                  "신규 화면 추가 시점에 즉시 회귀 검증 패턴 정착",
              ]),
         ]),
     ]),

    # === Chapter 8 — 종합 ===
    ("08", "RFP 요구사항 ↔ DataQ 한눈 정리",
     "DataQ 책임 영역 (A·C) + 차별화 4개",
     [
         ("커버리지 한 장 요약", [
             ("SFR-22 RFP 3개 영역 + DataQ 차별화는 어떻게 매핑되나요?",
              [
                  "**RFP A (메타 수집·관리) — DataQ 책임** → 다중 DBMS 자동 수집 + XMI 임포트 + 변경 이력 자동 기록",
                  "**RFP B (DB 모델링 관리) — ERwin 영역** → DataQ 는 ERwin 결과물을 임포트해 표준화·DDL 라인으로 연계",
                  "**RFP C (메타시스템 연계) — DataQ 책임** → Excel/DDL 추출·적재 + ERwin XMI 양방향 + REST API",
                  "**차별화 ①** 표준 사전 통합 + 자동 표준화 (DP + 동의어 cascade)",
                  "**차별화 ②** 데이터 품질 진단 — 룰 + 프로파일 + 시계열 통합",
                  "**차별화 ③** 거버넌스 + 운영 안전 — 권한 / 승인 / cascade / 인라인 편집 / dryRun",
                  "**차별화 ④** 자동화 — 진단 스케줄 + 약 100여 자동 회귀",
              ]),
         ]),
     ]),
]


# ========== 빌드 ==========
def build():
    # 슬라이드 카운트 사전 계산
    total = 1 + 1  # cover + overview
    for ch in CHAPTERS:
        total += 1   # section divider
        total += len(ch[3])  # content slides per chapter
    total += 1  # closing

    page_n = 1
    cover(); page_n += 1
    overview(total); page_n += 1
    for ch in CHAPTERS:
        idx, title, sub, contents = ch
        section(int(idx), title, sub, total, page_n); page_n += 1
        for header, pairs in contents:
            need_answer(page_n, total, header, pairs); page_n += 1
    closing(total)
    prs.save(OUT)
    print(f"saved: {OUT} ({total} slides)")


if __name__ == "__main__":
    build()
