"""
Narae DataQ 종합 시연 PPT 생성기 — 모던 디자인 v2.

스타일 컨셉:
  - Cover / Section / Closing: 다크 (#0E0E2C) 배경 + 그라디언트 액센트 도형
  - Content: 크림 (#FAFAFA) 배경 + 좌측 그라디언트 액센트 바 + 카드 레이아웃
  - 색상: violet(#7C3AED) + cyan(#06B6D4) 두 톤 액센트
  - 타이포: 큰 헤드라인 (44pt) + 미니멀 본문 (16pt) + 모노스페이스 코드
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Narae_DataQ_종합시연_2026-05-06.pptx")

# === 색상 팔레트 (Tailwind 기반) ===
DARK_BG    = RGBColor(0x0E, 0x0E, 0x2C)  # deep indigo near black
DARK_BG2   = RGBColor(0x1E, 0x1B, 0x4B)  # 지원
CREAM_BG   = RGBColor(0xFA, 0xFA, 0xFA)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK  = RGBColor(0x0F, 0x17, 0x2A)  # slate-900
TEXT_BODY  = RGBColor(0x33, 0x41, 0x55)  # slate-700
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)  # slate-500
TEXT_DIM   = RGBColor(0x94, 0xA3, 0xB8)  # slate-400
CARD_BORDER= RGBColor(0xE2, 0xE8, 0xF0)  # slate-200
DIVIDER    = RGBColor(0xCB, 0xD5, 0xE1)
# 액센트
VIOLET     = RGBColor(0x7C, 0x3A, 0xED)
INDIGO     = RGBColor(0x4F, 0x46, 0xE5)
CYAN       = RGBColor(0x06, 0xB6, 0xD4)
PINK       = RGBColor(0xEC, 0x48, 0x99)
# 의미 색
SUCCESS    = RGBColor(0x10, 0xB9, 0x81)
WARNING    = RGBColor(0xF5, 0x9E, 0x0B)
DANGER     = RGBColor(0xEF, 0x44, 0x44)
INFO_BLUE  = RGBColor(0x3B, 0x82, 0xF6)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

FONT_KR = '맑은 고딕'
FONT_MONO = 'Consolas'


# ─────────────────────────────────────────────────────────
# 헬퍼
# ─────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)


def fill_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(slide, left, top, width, height, color, line=False):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill_solid(s, color)
    if not line:
        s.line.fill.background()
    return s


def add_oval(slide, left, top, w, h, color, alpha=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, w, h)
    fill_solid(s, color)
    if alpha is not None:
        # alpha 적용 (XML)
        sp = s.fill.fore_color._xFill
        srgb = sp.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha_el = etree.SubElement(srgb, qn('a:alpha'))
            alpha_el.set('val', str(int(alpha * 100000)))
    return s


def add_text(slide, left, top, width, height, text, *, size=18, bold=False, color=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_KR, line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
    return box


def add_bg(slide, color=CREAM_BG):
    return add_rect(slide, 0, 0, SW, SH, color)


def add_decoration_dots(slide, base_color, n=20, area="full"):
    """배경 데코레이션 — 작은 원 + 큰 원 흩뿌리기 (alpha)"""
    import random
    random.seed(42)
    for _ in range(n):
        x = random.uniform(0, prs.slide_width.emu)
        y = random.uniform(0, prs.slide_height.emu)
        r = random.uniform(Inches(0.05).emu, Inches(0.25).emu)
        s = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(x), int(y), int(r), int(r))
        fill_solid(s, base_color)
        # alpha
        srgb = s.fill.fore_color._xFill.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha = etree.SubElement(srgb, qn('a:alpha'))
            alpha.set('val', '8000')  # 8% opacity


def add_accent_blob(slide, x, y, r, color, alpha=20000):
    """대형 블롭 (배경 그라디언트 효과 모방)"""
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, r, r)
    fill_solid(s, color)
    srgb = s.fill.fore_color._xFill.find(qn('a:srgbClr'))
    if srgb is not None:
        a = etree.SubElement(srgb, qn('a:alpha'))
        a.set('val', str(alpha))


def add_left_accent_bar(slide):
    """좌측 5px 액센트 바 (violet → cyan 두 단)"""
    add_rect(slide, 0, 0, Inches(0.12), Inches(3.75), VIOLET)
    add_rect(slide, 0, Inches(3.75), Inches(0.12), Inches(3.75), CYAN)


def add_section_chip(slide, left, top, label, color=VIOLET):
    """SECTION · 4-2 같은 칩"""
    box = slide.shapes.add_textbox(left, top, Inches(3.0), Inches(0.4))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.name = FONT_KR
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = color
    return box


def add_card(slide, left, top, width, height, *, fill=WHITE, border=CARD_BORDER):
    """카드 (라운드 사각형) — 그림자 대신 얇은 테두리"""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.adjustments[0] = 0.06  # corner radius
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = border
    s.line.width = Pt(0.75)
    return s


def add_bullet_list(slide, left, top, width, height, items, *, size=15,
                    head_color=TEXT_DARK, body_color=TEXT_BODY, marker_color=VIOLET):
    """
    items: ["헤딩", (1, "본문"), (1, "본문2"), "헤딩2", ...]
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_top = Emu(0)
    first = True
    for item in items:
        if isinstance(item, tuple):
            indent_lvl, txt = item
        else:
            indent_lvl, txt = 0, item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.25
        if indent_lvl == 0:
            run = p.add_run()
            run.text = "▍ " + txt
            run.font.name = FONT_KR
            run.font.size = Pt(size + 1)
            run.font.bold = True
            run.font.color.rgb = head_color
        else:
            indent = "    "
            r1 = p.add_run()
            r1.text = indent + "·  "
            r1.font.name = FONT_KR
            r1.font.size = Pt(size)
            r1.font.color.rgb = marker_color
            r1.font.bold = True
            r2 = p.add_run()
            r2.text = txt
            r2.font.name = FONT_KR
            r2.font.size = Pt(size)
            r2.font.color.rgb = body_color
    return box


def add_page_number(slide, n, total=36, color=TEXT_DIM):
    add_text(slide, Inches(12.4), Inches(7.0), Inches(0.8), Inches(0.4),
             f"{n:02d} / {total}", size=10, color=color, align=PP_ALIGN.RIGHT, font='Consolas')


def add_brand_footer(slide, color=TEXT_DIM):
    add_text(slide, Inches(0.4), Inches(7.0), Inches(4.0), Inches(0.4),
             "Narae DataQ · 종합 시연 2026-05-06", size=10, color=color, font=FONT_KR)


def add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def gradient_band(slide, left, top, width, height, c1, c2, vertical=False):
    """단순 그라디언트 — 가로로 N개 사각형 분할"""
    n = 50
    for i in range(n):
        # 보간
        t = i / (n - 1)
        r = int(c1[0] + (c2[0] - c1[0]) * t)
        g = int(c1[1] + (c2[1] - c1[1]) * t)
        b = int(c1[2] + (c2[2] - c1[2]) * t)
        if vertical:
            seg_h = int(height / n)
            seg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         int(left), int(top + i * seg_h),
                                         int(width), int(seg_h + 1))
        else:
            seg_w = int(width / n)
            seg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         int(left + i * seg_w), int(top),
                                         int(seg_w + 1), int(height))
        seg.fill.solid()
        seg.fill.fore_color.rgb = RGBColor(r, g, b)
        seg.line.fill.background()


def gradient_xml(shape, c1, c2, angle=2700000):
    """그라디언트 XML 직접 삽입 (angle: 1/60000 deg, 0=오른쪽, 5400000=아래)"""
    spPr = shape.fill._xPr
    # 기존 fill 제거
    for tag in ['a:solidFill', 'a:gradFill', 'a:noFill', 'a:blipFill', 'a:pattFill']:
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    grad = etree.SubElement(spPr, qn('a:gradFill'))
    grad.set('rotWithShape', '1')
    gsLst = etree.SubElement(grad, qn('a:gsLst'))
    for pos, c in [(0, c1), (100000, c2)]:
        gs = etree.SubElement(gsLst, qn('a:gs'))
        gs.set('pos', str(pos))
        srgb = etree.SubElement(gs, qn('a:srgbClr'))
        srgb.set('val', f'{c[0]:02X}{c[1]:02X}{c[2]:02X}')
    lin = etree.SubElement(grad, qn('a:lin'))
    lin.set('ang', str(angle))
    lin.set('scaled', '0')


# ─────────────────────────────────────────────────────────
# 슬라이드 빌더
# ─────────────────────────────────────────────────────────

def title_slide(num, label, title, subtitle, hint):
    """다크 표지/섹션 슬라이드"""
    s = add_slide()
    add_bg(s, DARK_BG)
    # 큰 그라디언트 블롭 2개 (장식)
    add_accent_blob(s, Inches(8.5), Inches(-2.0), Inches(7.0).emu, VIOLET, alpha=20000)
    add_accent_blob(s, Inches(-2.0), Inches(4.5), Inches(6.0).emu, CYAN, alpha=15000)
    add_accent_blob(s, Inches(10.5), Inches(5.0), Inches(4.0).emu, PINK, alpha=10000)
    # 좌상단 라벨
    add_text(s, Inches(0.7), Inches(0.6), Inches(6.0), Inches(0.4),
             label, size=12, bold=True, color=CYAN, font='Consolas')
    # 메인 타이틀
    add_text(s, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.5),
             title, size=64, bold=True, color=WHITE)
    # 서브
    if subtitle:
        add_text(s, Inches(0.7), Inches(3.8), Inches(12.0), Inches(0.8),
                 subtitle, size=24, color=RGBColor(0xC7, 0xD2, 0xFE))
    # 액센트 라인
    add_rect(s, Inches(0.7), Inches(4.7), Inches(0.8), Pt(4), VIOLET)
    add_rect(s, Inches(1.6), Inches(4.7), Inches(0.5), Pt(4), CYAN)
    # 힌트
    if hint:
        add_text(s, Inches(0.7), Inches(5.2), Inches(12.0), Inches(1.5),
                 hint, size=14, color=TEXT_DIM)
    # 페이지 번호
    add_text(s, Inches(12.4), Inches(7.0), Inches(0.8), Inches(0.4),
             f"{num:02d} / 36", size=10, color=TEXT_DIM, align=PP_ALIGN.RIGHT, font='Consolas')
    add_text(s, Inches(0.4), Inches(7.0), Inches(4.0), Inches(0.4),
             "Narae DataQ", size=10, color=TEXT_DIM, font=FONT_KR)
    return s


def content_slide(num, chip, title, subtitle, items, *, accent=VIOLET, accent2=CYAN):
    """라이트 콘텐츠 슬라이드 — 좌측 액센트 바 + 카드 레이아웃"""
    s = add_slide()
    add_bg(s, CREAM_BG)
    # 좌측 그라디언트 바
    gradient_band(s, 0, 0, Inches(0.18), SH, (accent[0], accent[1], accent[2]),
                  (accent2[0], accent2[1], accent2[2]), vertical=True)
    # 칩
    add_section_chip(s, Inches(0.7), Inches(0.55), chip, color=accent)
    # 제목
    add_text(s, Inches(0.7), Inches(0.95), Inches(12.0), Inches(0.85),
             title, size=32, bold=True, color=TEXT_DARK)
    # 서브
    if subtitle:
        add_text(s, Inches(0.7), Inches(1.85), Inches(12.0), Inches(0.5),
                 subtitle, size=15, color=TEXT_MUTED)
    # 디바이더 (얇은 선)
    add_rect(s, Inches(0.7), Inches(2.45), Inches(12.0), Pt(0.5), DIVIDER)
    # 콘텐츠 카드 (단일)
    card = add_card(s, Inches(0.7), Inches(2.7), Inches(12.0), Inches(4.1),
                    fill=WHITE, border=CARD_BORDER)
    add_bullet_list(s, Inches(1.0), Inches(2.95), Inches(11.4), Inches(3.7),
                    items, size=15, marker_color=accent)
    # 페이지 / 브랜드
    add_page_number(s, num)
    add_brand_footer(s)
    return s


def two_col_slide(num, chip, title, subtitle, left_items, right_items, *,
                  left_label=None, right_label=None, accent=VIOLET, accent2=CYAN):
    """2단 콘텐츠"""
    s = add_slide()
    add_bg(s, CREAM_BG)
    gradient_band(s, 0, 0, Inches(0.18), SH, (accent[0], accent[1], accent[2]),
                  (accent2[0], accent2[1], accent2[2]), vertical=True)
    add_section_chip(s, Inches(0.7), Inches(0.55), chip, color=accent)
    add_text(s, Inches(0.7), Inches(0.95), Inches(12.0), Inches(0.85),
             title, size=32, bold=True, color=TEXT_DARK)
    if subtitle:
        add_text(s, Inches(0.7), Inches(1.85), Inches(12.0), Inches(0.5),
                 subtitle, size=15, color=TEXT_MUTED)
    add_rect(s, Inches(0.7), Inches(2.45), Inches(12.0), Pt(0.5), DIVIDER)
    # 좌측 카드
    add_card(s, Inches(0.7), Inches(2.7), Inches(5.95), Inches(4.1))
    if left_label:
        add_text(s, Inches(0.95), Inches(2.85), Inches(5.5), Inches(0.4),
                 left_label, size=11, bold=True, color=accent, font='Consolas')
    add_bullet_list(s, Inches(0.95), Inches(3.25), Inches(5.5), Inches(3.4),
                    left_items, size=14, marker_color=accent)
    # 우측 카드
    add_card(s, Inches(6.75), Inches(2.7), Inches(5.95), Inches(4.1))
    if right_label:
        add_text(s, Inches(7.0), Inches(2.85), Inches(5.5), Inches(0.4),
                 right_label, size=11, bold=True, color=accent2, font='Consolas')
    add_bullet_list(s, Inches(7.0), Inches(3.25), Inches(5.5), Inches(3.4),
                    right_items, size=14, marker_color=accent2)
    add_page_number(s, num)
    add_brand_footer(s)
    return s


def stat_slide(num, chip, title, subtitle, stats, items=None, *, accent=VIOLET, accent2=CYAN):
    """통계 KPI 카드 + 본문"""
    s = add_slide()
    add_bg(s, CREAM_BG)
    gradient_band(s, 0, 0, Inches(0.18), SH, (accent[0], accent[1], accent[2]),
                  (accent2[0], accent2[1], accent2[2]), vertical=True)
    add_section_chip(s, Inches(0.7), Inches(0.55), chip, color=accent)
    add_text(s, Inches(0.7), Inches(0.95), Inches(12.0), Inches(0.85),
             title, size=32, bold=True, color=TEXT_DARK)
    if subtitle:
        add_text(s, Inches(0.7), Inches(1.85), Inches(12.0), Inches(0.5),
                 subtitle, size=15, color=TEXT_MUTED)
    add_rect(s, Inches(0.7), Inches(2.45), Inches(12.0), Pt(0.5), DIVIDER)
    # KPI 카드 (3~6개)
    n = len(stats)
    card_w = (Inches(12.0) - Inches(0.2) * (n - 1)) / n
    for i, st in enumerate(stats):
        x = Inches(0.7) + i * (card_w + Inches(0.2))
        add_card(s, x, Inches(2.7), card_w, Inches(1.4), fill=WHITE)
        # 큰 숫자
        add_text(s, x, Inches(2.85), card_w, Inches(0.7),
                 st['value'], size=32, bold=True, color=st.get('color', accent),
                 align=PP_ALIGN.CENTER)
        # 라벨
        add_text(s, x, Inches(3.55), card_w, Inches(0.4),
                 st['label'], size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    # 본문 카드
    if items:
        add_card(s, Inches(0.7), Inches(4.3), Inches(12.0), Inches(2.5))
        add_bullet_list(s, Inches(1.0), Inches(4.5), Inches(11.4), Inches(2.2),
                        items, size=14, marker_color=accent)
    add_page_number(s, num)
    add_brand_footer(s)
    return s


# ─────────────────────────────────────────────────────────
# 슬라이드 생성
# ─────────────────────────────────────────────────────────

# ━━━ 1. 표지 ━━━
s = add_slide()
add_bg(s, DARK_BG)
# 큰 블롭들
add_accent_blob(s, Inches(8.0), Inches(-2.5), Inches(8.0).emu, VIOLET, alpha=25000)
add_accent_blob(s, Inches(-3.0), Inches(3.5), Inches(7.0).emu, CYAN, alpha=18000)
add_accent_blob(s, Inches(9.5), Inches(4.5), Inches(5.0).emu, PINK, alpha=12000)
# 라벨
add_text(s, Inches(0.7), Inches(0.6), Inches(6.0), Inches(0.4),
         "DEMO · 2026-05-06", size=12, bold=True, color=CYAN, font='Consolas')
# 메인
add_text(s, Inches(0.7), Inches(2.0), Inches(12.0), Inches(1.6),
         "Narae DataQ", size=84, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(3.7), Inches(12.0), Inches(0.7),
         "데이터 품질 · 표준화 관리 플랫폼", size=28, color=RGBColor(0xC7, 0xD2, 0xFE))
# 액센트 바
add_rect(s, Inches(0.7), Inches(4.65), Inches(0.8), Pt(5), VIOLET)
add_rect(s, Inches(1.6), Inches(4.65), Inches(0.5), Pt(5), CYAN)
add_rect(s, Inches(2.2), Inches(4.65), Inches(0.3), Pt(5), PINK)
# 서브 텍스트
add_text(s, Inches(0.7), Inches(5.0), Inches(12.0), Inches(1.5),
         ["종합 시연  —  단어 · 용어 · 도메인  →  데이터 모델  →  표준 / 구조 / 품질 진단",
          "자동 표준화  ·  진단 스케줄  ·  진단 제외 관리"],
         size=14, color=TEXT_DIM, line_spacing=1.6)
add_text(s, Inches(0.7), Inches(6.85), Inches(6.0), Inches(0.4),
         "발표  —  장재영", size=12, color=TEXT_DIM, font=FONT_KR)
add_text(s, Inches(12.4), Inches(7.0), Inches(0.8), Inches(0.4),
         "01 / 36", size=10, color=TEXT_DIM, align=PP_ALIGN.RIGHT, font='Consolas')
add_notes(s, "안녕하세요. Narae DataQ — 데이터 품질·표준화 관리 플랫폼의 종합 시연을 시작하겠습니다. "
            "단어·용어·도메인 표준 관리부터 데이터 모델 수집·표준화 진단·구조 변경 감지·"
            "데이터 품질 검증·진단 스케줄링까지, 데이터 표준화 라이프사이클 전체를 한 솔루션에서 다룹니다. "
            "이번 시연은 좌측 메뉴 순서대로 모든 핵심 기능을 보여드리는 종합 시연입니다.")


# ━━━ 2. Agenda ━━━
s = add_slide()
add_bg(s, CREAM_BG)
gradient_band(s, 0, 0, Inches(0.18), SH, (VIOLET[0], VIOLET[1], VIOLET[2]),
              (CYAN[0], CYAN[1], CYAN[2]), vertical=True)
add_section_chip(s, Inches(0.7), Inches(0.55), "AGENDA")
add_text(s, Inches(0.7), Inches(0.95), Inches(12.0), Inches(0.85),
         "Agenda", size=36, bold=True, color=TEXT_DARK)
add_text(s, Inches(0.7), Inches(1.85), Inches(12.0), Inches(0.5),
         "메뉴 순서를 그대로 따라가는 종합 시연 — 12 영역 / 36 슬라이드 / 50분",
         size=15, color=TEXT_MUTED)
add_rect(s, Inches(0.7), Inches(2.45), Inches(12.0), Pt(0.5), DIVIDER)

# 12개 카드 grid (4 cols × 3 rows)
agenda = [
    ("01", "시스템 개요", "Java + Spring Boot + Vue + Postgres", VIOLET),
    ("02", "로그인 + 대시보드", "권한 + KPI 카드", VIOLET),
    ("03", "데이터 표준 사전", "단어 / 용어 v2 / 코드 / 도메인", INDIGO),
    ("04", "데이터 모델", "수집 / 그리드 / DDL / 진단 제외", INDIGO),
    ("05", "표준화 진단", "실행 / 결과 / 빠른 등록", CYAN),
    ("06", "구조 변경 진단", "스키마 vs 스냅샷 비교", CYAN),
    ("07", "자동 표준화 지원", "DP + 동의어 cascade", PINK),
    ("08", "데이터 품질 진단", "룰 / 위반 샘플 / 시계열", PINK),
    ("09", "진단 스케줄", "cron + 권한 + 동시 방어", VIOLET),
    ("10", "마이페이지 + 관리", "요청 현황 / 사용자 / 데이터소스", VIOLET),
    ("11", "셀레니움 회귀", "29 PASS / 0 FAIL / 1675초", SUCCESS),
    ("12", "SMETA 비교 + Q&A", "차별화 4 영역", PINK),
]
gx, gy = Inches(0.7), Inches(2.7)
cw, ch = Inches(2.95), Inches(1.35)
gap = Inches(0.07)
for i, (num, t, sub, c) in enumerate(agenda):
    row = i // 4; col = i % 4
    x = gx + col * (cw + gap); y = gy + row * (ch + gap)
    add_card(s, x, y, cw, ch, fill=WHITE)
    # 번호 (액센트)
    add_text(s, x + Inches(0.2), y + Inches(0.13), Inches(0.6), Inches(0.3),
             num, size=10, bold=True, color=c, font='Consolas')
    # 제목
    add_text(s, x + Inches(0.2), y + Inches(0.4), cw - Inches(0.4), Inches(0.4),
             t, size=14, bold=True, color=TEXT_DARK)
    # 서브
    add_text(s, x + Inches(0.2), y + Inches(0.8), cw - Inches(0.4), Inches(0.4),
             sub, size=10, color=TEXT_MUTED)
    # 좌측 액센트 바
    add_rect(s, x, y + Inches(0.13), Inches(0.04), Inches(1.05), c)

add_page_number(s, 2)
add_brand_footer(s)
add_notes(s, "오늘 시연은 좌측 네비게이션 메뉴 순서를 그대로 따라가면서 12개 영역의 핵심 기능을 보여드립니다. "
            "마지막에 셀레니움 자동 회귀 결과와 SMETA 대비 차별화 포인트를 정리합니다.")


# ━━━ 3. 환경 ━━━
content_slide(3, "01  ·  ENVIRONMENT", "시스템 환경",
              "Java 1.8  +  Spring Boot 2.7  +  Vue 2  +  PostgreSQL 13",
              [
                  "Backend (2 modules)",
                  (1, "q-center  —  웹 서버 / 28091 / Controller + Vue 프론트"),
                  (1, "q-executor  —  백그라운드 워커 / 28098 / 진단·스케줄 실행"),
                  (1, "q-common  —  공유 VO + MyBatis Mapper"),
                  "Frontend",
                  (1, "Vue 2.5 + Vuetify 2.6 (SPA, keep-alive 탭)"),
                  (1, "WebSocket (STOMP)  —  진단 실시간 진행률 push"),
                  "Database / Build / Deploy",
                  (1, "메타: PostgreSQL 13 (외부 25433)  ·  외부: Oracle SID/Service, Cubrid, MySQL/MariaDB"),
                  (1, "Maven 루트 reactor (~2분)  ·  Docker 컨테이너 (dataq-db, oracle-xe)"),
                  (1, "DDL_full_schema.sql 단일 진실  →  신규 환경 1 파일 실행"),
              ])
add_notes(prs.slides[2], "백엔드는 Spring Boot 두 모듈 — q-center 웹과 q-executor 워커. 프론트는 Vue 2 SPA, "
            "데이터 모델 진단은 WebSocket 으로 실시간 진행률 push. 메타DB 는 PostgreSQL 13 컨테이너이고, "
            "외부 데이터소스로 Oracle SID/Service Name 둘 다 지원, Cubrid 호환성도 확보했습니다. "
            "DDL 은 pg_dump 결과 한 파일로 단일 진실 운영 — 신규 환경 구축은 그 파일 한 번 실행으로 끝납니다.")


# ━━━ 4. 로그인 ━━━
content_slide(4, "02  ·  LOGIN", "로그인",
              "ID/PW + SHA-256 솔트 + 권한 (admin / user 2단계)",
              [
                  "시연 절차",
                  (1, "http://localhost:28091  →  ID 'space' (관리자) / PW '123'"),
                  (1, "[로그인]  →  /app/main 진입"),
                  "검증 포인트",
                  (1, "비밀번호: SHA-256 + 솔트 (ndata-quality-secret)"),
                  (1, "세션: HttpSession + STOMP WebSocket 동시 인증"),
                  (1, "권한: TB_USER.ADMIN_YN — 'Y' 면 [관리] 메뉴 노출"),
                  "Fallback",
                  (1, "DB 연결 실패 시  →  docker restart dataq-db  →  재시도"),
              ])
add_notes(prs.slides[3], "ID 와 비밀번호는 SHA-256 + 솔트로 저장합니다. 권한은 admin/user 두 단계인데, "
            "중요한 건 admin 만 [관리] 그룹 메뉴가 DOM 에 렌더되고 API 도 403 거부합니다. 지금 로그인합니다 — 'space' / '123'.")


# ━━━ 5. 대시보드 ━━━
content_slide(5, "03  ·  DASHBOARD", "대시보드",
              "표준 현황 + 데이터 모델 현황 + 승인 현황 + 추이 차트",
              [
                  "카드 4개 (클릭 시 해당 메뉴로 이동)",
                  (1, "표준 현황  —  단어/용어/도메인 승인된 건수"),
                  (1, "데이터 모델 현황  —  모델 / 테이블 / 컬럼 합계"),
                  (1, "승인 현황  —  대기 / 승인 완료 / 반려"),
                  (1, "표준화 준수율  —  모델별 percent"),
                  "추이 차트",
                  (1, "최근 N회 진단의 표준 준수율 시계열"),
                  (1, "준수율 = (전체 컬럼수 − 이슈 컬럼수) / 전체 컬럼수 × 100"),
                  (1, "RESULT_CNT(이슈 건수) 가 아닌 ISSUE_COL_CNT(이슈 컬럼수) 사용"),
              ], accent=INDIGO)
add_notes(prs.slides[4], "대시보드는 한 화면 안에 표준화 운영 KPI 를 보여줍니다. 카드 클릭하면 해당 메뉴로 점프. "
            "표준화 준수율은 컬럼 단위 — DISTINCT(테이블.컬럼) 로 카운트해서 한 컬럼이 여러 이슈여도 1로 처리. "
            "이렇게 해야 의미 있는 백분율이 나옵니다.")


# ━━━ 6. 단어 ━━━
content_slide(6, "04-1  ·  WORD", "데이터 표준 사전 — 단어",
              "행안부 공통 표준 단어 + 사용자 추가 + 형식단어 (분류어)",
              [
                  "검색",
                  (1, "단어명 / 영문약어 / 형식단어 여부 / 등록일자 범위 / 승인 여부"),
                  "그리드",
                  (1, "행 클릭 → 상세 패널 (영문 풀명 / 단어 설명 / 동의어 리스트)"),
                  "버튼",
                  (1, "[등록] / [일괄 등록] / [템플릿 다운로드] / [다운로드] / [삭제]"),
                  "등록 모달",
                  (1, "한글 단어명 → 영문약어 (자동 대문자) → 영문 풀명"),
                  (1, "도메인 분류 + 형식단어 여부 토글  ·  동의어 리스트 + 금칙어 리스트"),
              ], accent=INDIGO)
add_notes(prs.slides[5], "단어 사전은 행안부 공통 표준 약 3,300건 + 우리 추가 단어. 검색 필터에 등록일자 범위가 있어서 "
            "특정 기간 내 등록된 단어만 추적할 수 있고, 형식단어 여부로 명사·분류어 별도 검색이 가능합니다. "
            "등록 모달에서 동의어 리스트가 핵심 — 자동 표준화 시 이 동의어를 통해 미등록 단어를 자동 매칭합니다.")


# ━━━ 7. 용어 v2 ━━━
content_slide(7, "04-2  ·  TERM v2  (NEW)", "용어 등록 v2 — 81/82번",
              "단일 폼  +  1초 디바운스 자동 분석  +  코드 picker",
              [
                  "v1 → v2 핵심 변화",
                  (1, "3-step stepper 폐지  →  단일 폼"),
                  (1, "한글 용어명 → 1초 디바운스 → analyzeTermsBatch API 자동 호출"),
                  (1, "DSTermRecommend 와 동일 알고리즘 (DP 점수 + 동의어 cascade)"),
                  "응답 — 단어 1개당 분류 1개 (부분문자열 잡음 제거)",
                  (1, "MATCHED  →  wordLst[selected wordCandidate] + 자동 체크"),
                  (1, "NEW / UNRECOGNIZED  →  인라인 등록 폼 즉시 노출"),
                  "추가 기능",
                  (1, "추천 도메인 자동 채움 (recommendedDomainNm)"),
                  (1, "마지막 단어가 'CD'  →  도메인 유형 자동 '코드' 토글 + picker"),
              ], accent=PINK)
add_notes(prs.slides[6], "v1 은 stepper 3 단계로 사용자 클릭이 많았고 부분문자열 매칭이 너무 많이 노출됐습니다. "
            "예를 들어 '가로세로일시' 입력 시 27개 분류가 동시에 나와서 어디를 선택해야 할지 헷갈렸어요. "
            "v2 는 백엔드의 analyzeTermsBatch API 로 교체했는데, 이게 DP 점수 + 가장 긴 매칭 우선 + 동의어 cascade 적용해서 "
            "단어 1개당 분류 1개만 깨끗하게 반환합니다. MATCHED 단어는 자동 체크되고, 추천 도메인은 자동 채움. "
            "마지막 단어가 'CD' 면 코드 picker 가 자동으로 활성화. 사용자 클릭 횟수가 v1 대비 3분의 1 수준으로 줄었습니다.")


# ━━━ 8. 용어 분석 4종 ━━━
two_col_slide(8, "04-2  ·  TERM v2  ·  CASES", "용어 분석 4종 status",
              "AUTO / PARTIAL / FAILED / REGISTERED  ·  자동 셀레니움 검증 8/8 PASS",
              left_items=[
                  "AUTO  —  모든 단어 매칭",
                  (1, "예: '회원전화번호'  →  [회원, 전화, 번호]"),
                  (1, "→ MBR_TEL_NO 즉시 등록 가능"),
                  "PARTIAL  —  일부 미등록",
                  (1, "예: '블라블라일자'  →  '일자' MATCHED + '블라블라' UNRECOGNIZED"),
                  (1, "→ 미등록 칩 + 인라인 등록 폼"),
              ],
              right_items=[
                  "FAILED  —  전 토큰 미인식",
                  (1, "예: '라랄라룰루'  →  모든 토큰 UNRECOGNIZED"),
                  (1, "→ swal 안내"),
                  "REGISTERED  —  이미 등록된 용어",
                  (1, "중복 swal 알림"),
                  (1, "→ 자동 분석은 그대로 표시",),
                  "→ test_term_register_v2.py 8/8 PASS",
              ],
              left_label="MATCHED",
              right_label="NEW / DUP",
              accent=PINK, accent2=VIOLET)
add_notes(prs.slides[7], "용어 분석 결과는 4가지로 분류됩니다. AUTO 는 모두 매칭 — 그대로 등록 가능. "
            "PARTIAL 은 일부 미등록인데, 미등록 단어를 모달 안에서 인라인으로 즉시 등록할 수 있어서 화면 전환 없이 흐름이 끊기지 않습니다. "
            "FAILED 는 다시 입력하라는 안내, REGISTERED 는 이미 있으니 중복이라고 알려줍니다. "
            "이 4가지가 셀레니움 8 케이스로 자동 검증됩니다.")


# ━━━ 9. 코드/도메인/그룹/분류/이력 ━━━
content_slide(9, "04-3  ·  CODE  ·  DOMAIN  ·  HISTORY",
              "코드 / 도메인 / 그룹 / 분류 / 변경 이력",
              "표준 사전의 나머지 5개 메뉴",
              [
                  "코드  —  코드 그룹 + 코드 항목",
                  (1, "예: '성별 (GENDER_CD)' = M / F / N  ·  등록일자 범위 검색 신규"),
                  "도메인  —  데이터 타입/길이 표준",
                  (1, "예: '금액 (AMT)' = NUMBER(15,2)"),
                  "도메인 그룹 / 도메인 분류",
                  (1, "도메인을 분류로 묶어 자동 표준화 시 cascade 추천"),
                  "변경 이력",
                  (1, "단어/용어/도메인/코드 모든 항목 등록·수정·삭제 이력"),
                  (1, "변경 사용자, 시각, before/after 값 모두 추적"),
              ], accent=INDIGO)
add_notes(prs.slides[8], "코드는 코드 그룹과 항목 — 예를 들어 성별 코드 그룹에 M, F, N 같은 항목들. "
            "도메인은 데이터 타입+길이 표준화. 도메인을 분류로 묶으면 자동 표준화 시 마지막 단어의 분류명에 따라 도메인이 cascade 됩니다. "
            "변경 이력은 모든 항목의 before/after 값을 보존 — 누가 언제 뭘 바꿨는지 100% 추적 가능합니다.")


# ━━━ 10. 일괄 등록 ━━━
content_slide(10, "04-4  ·  BULK  ·  TEMPLATES", "일괄 등록 양식 + 결과 다운로드",
              "정적 XLSX 5종  +  동적 POI",
              [
                  "정적 XLSX 양식 (5종)",
                  (1, "단어 / 용어 / 도메인 / 테이블 / 컬럼"),
                  (1, "[템플릿 다운로드]  →  즉시 다운로드 (ClassPathResource)"),
                  "일괄 등록 흐름",
                  (1, "양식 입력  →  [일괄 등록]  →  미리보기 (검증 결과 색 구분)  →  [커밋]"),
                  (1, "검증 실패 행은 빨강 + 사유 인라인"),
                  "결과 다운로드 (동적)",
                  (1, "검색 필터 적용 결과를 POI 동적 생성으로 XLSX"),
                  (1, "현재 그리드 상태 그대로 — 사용자 의도 반영"),
              ], accent=INDIGO)
add_notes(prs.slides[9], "양식 다운로드는 정적 XLSX 5종 — 미리 디자인된 헤더만 있는 빈 양식. 사용자가 입력해서 다시 일괄 등록하면 미리보기 단계에서 "
            "검증 결과를 색으로 구분해서 보여주고 사용자 확인 후 커밋. 결과 다운로드는 정적이 아니라 POI 로 동적 생성 — "
            "현재 검색·필터 상태 그대로 다운로드되니까 사용자 의도가 반영됩니다.")


# ━━━ 11. 승인 워크플로 ━━━
content_slide(11, "04-5  ·  APPROVAL", "승인 워크플로 (관리 메뉴)",
              "단어 선승인  +  cascade 반려  +  반려 후 물리 삭제",
              [
                  "승인 흐름",
                  (1, "일반 사용자 등록 → APRV_YN='N' → [관리 > 승인] 대기열"),
                  (1, "관리자 행별 [승인]/[반려] 인라인 (사유 입력)"),
                  (1, "관리자 등록 시 APRV_YN='Y' 즉시 승인"),
                  "정책",
                  (1, "단어 선승인 — 구성 단어 미승인 시 용어 승인 거부 (alert)"),
                  (1, "cascade 반려 — 단어 반려 시 연관 미승인 용어 동시 반려 + 알림"),
                  (1, "반려 후 물리 삭제 — 동일 단어명 재등록 가능"),
                  "검증",
                  (1, "test_word_approval_flow / test_full_approval_flow / test_cascade_and_word_first / test_reject_physical_delete"),
              ], accent=INDIGO)
add_notes(prs.slides[10], "승인 워크플로는 일반 사용자 등록 시 APRV_YN='N' 으로 들어가고 관리자가 승인 화면에서 행별로 승인/반려. "
            "정책이 핵심입니다 — 단어가 승인 안 되면 그 단어를 쓴 용어도 승인 거부. 단어 반려하면 연관된 미승인 용어도 동시 반려. "
            "반려된 항목은 물리 삭제해서 동일 이름 재등록이 가능. 이 정책들이 셀레니움 4건 + cascade 12건으로 자동 검증됩니다.")


# ━━━ 12. 데이터 모델 등록 ━━━
content_slide(12, "05-1  ·  DATA MODEL", "데이터 모델 — 등록 + 자동 수집",
              "Oracle / Postgres / Cubrid / MySQL / MariaDB  메타데이터 자동 수집",
              [
                  "데이터 소스 등록 (관리 > 데이터 소스)",
                  (1, "DBMS 유형  —  Oracle SID / Service Name / Postgres / Cubrid / MySQL / MariaDB"),
                  (1, "호스트 / 포트 / 계정 / 비밀번호 (jasypt 암호화)  ·  [연결 테스트]"),
                  "데이터 모델 등록 (데이터 모델 > 관리)",
                  (1, "데이터소스 + 스키마 다중 선택 → 저장 + 자동 수집"),
                  (1, "수집: 테이블 + 컬럼 + 인덱스 + 제약조건 4 영역"),
                  (1, "재수집 시: ADDED_CNT / DELETED_CNT / MODIFIED_CNT 통계 (44/48번)"),
                  "물리 모델 / 논리 모델 분리",
                  (1, "DSID NULL 이면 논리 모델  →  표준 진단만 (구조 진단 거부)"),
              ], accent=CYAN)
add_notes(prs.slides[11], "외부 DBMS 메타데이터를 자동으로 수집합니다. Oracle 은 SID 와 Service Name 둘 다 지원. "
            "Cubrid 는 시스템 카탈로그 호환성을 추가로 작업해서 사용 가능. 모델 등록 시 스키마 다중 선택해서 한 번에 수집. "
            "재수집 시 추가/삭제/변경 통계가 자동 누적되는데, 이 컬럼이 어제 PC1 과 동기화 작업한 부분입니다. "
            "물리 모델은 데이터소스 연결돼있고, 논리 모델은 DSID 없는 모델 — 둘이 진단 적용 범위가 달라요.")


# ━━━ 13. 수집 이력 ━━━
content_slide(13, "05-1-2  ·  COLLECT LOG", "수집 이력",
              "수집 이벤트 로그 + 변경 통계",
              [
                  "테이블 컬럼",
                  (1, "수집 ID / 모델명 / 시작·종료 시각 / 완료 여부"),
                  (1, "ADDED_CNT  —  신규 추가된 테이블·컬럼 합계"),
                  (1, "DELETED_CNT  —  삭제된 테이블·컬럼 합계"),
                  (1, "MODIFIED_CNT  —  변경된 테이블·컬럼 합계"),
                  "활용",
                  (1, "운영 추적  —  '모델별 마지막 수집 + 변경량' 한 화면 조회"),
                  (1, "장애 시  —  '어제 vs 오늘 수집' 차이로 원인 추정"),
                  (1, "감사  —  누가 언제 수집을 트리거했는지 변경 이력에 자동 등록"),
              ], accent=CYAN)
add_notes(prs.slides[12], "수집 이력은 단순한 로그가 아니라 변경 통계까지 같이 표시되는 감사 로그입니다. "
            "44번 설계에서 정의한 ADDED_CNT/DELETED_CNT/MODIFIED_CNT 가 핵심 — 운영 중에 어떤 모델에 변경이 컸는지 한눈에 보입니다.")


# ━━━ 14. 그리드 편집 ━━━
content_slide(14, "05-2  ·  GRID EDIT  (53번)", "테이블 + 컬럼 그리드 편집",
              "3가지 입력 경로  +  인라인 검증  +  dirty 표시",
              [
                  "3가지 입력 경로",
                  (1, "그리드 직접  —  인라인 행 추가, 셀 편집, 한글명 중복 검증"),
                  (1, "TSV 붙여넣기  —  엑셀 복사 → Ctrl+V → 자동 행 분배"),
                  (1, "엑셀 일괄 업로드  —  양식 다운로드 → 입력 → 미리보기 → 커밋"),
                  "헤더 라벨 통일 (PC1 4-27 보강)",
                  (1, "'테이블명' → '테이블 영문명 (물리)' 등 명확화"),
                  (1, "헤더 3색  ·  인라인 편집  ·  dirty 노란 배경"),
                  "변환 (자동 표준화 연동)",
                  (1, "한글명 입력 → [변환] → 영문약어 / 타입 / 길이 자동"),
                  (1, "변환 실패 시 → TMP_COL_{n} + VARCHAR(255) 비표준 자동 저장"),
              ], accent=CYAN)
add_notes(prs.slides[13], "컬럼 편집은 53번 설계로 재작업한 영역인데, 입력 경로가 3개입니다. "
            "그리드 직접 편집, TSV 붙여넣기, 엑셀 업로드. TSV 붙여넣기가 운영 중 가장 자주 쓰이는데 "
            "엑셀에서 표 복사해서 그리드에 붙이면 행이 자동으로 분배됩니다. "
            "한글명을 입력하고 [변환] 누르면 자동 표준화 엔진이 영문약어와 타입을 자동으로 채워줍니다.")


# ━━━ 15. cascade rename ━━━
content_slide(15, "05-2-2  ·  CASCADE RENAME", "테이블 물리명 변경 — cascade",
              "5단계 cascade  ·  OBJ_NM / ATTR / INDEX / CONSTRAINT / REF_TABLE_NM",
              [
                  "변경 영향도 미리보기",
                  (1, "previewObjRename API  —  5개 카운트"),
                  (1, "(OBJ + ATTR + INDEX + CONSTRAINT + REF_TABLE)"),
                  (1, "사용자 확인 swal  →  진행 여부 결정"),
                  "실제 변경",
                  (1, "updateObj API  —  5단계 cascade rename 트랜잭션"),
                  (1, "DM_ID + 자연키 PK 기준 (CLCT 폐기 후)"),
                  "검증",
                  (1, "test_obj_rename_cascade.py — 8/8 PASS"),
              ], accent=CYAN)
add_notes(prs.slides[14], "테이블 물리명을 바꿀 때 그 테이블 컬럼들, 인덱스, 제약조건, 그리고 외래키로 이 테이블을 참조하는 다른 테이블의 REF_TABLE_NM 까지 "
            "5단계로 cascade rename 됩니다. 변경 전에 영향도 카운트가 미리보기로 swal 에 뜨고 사용자 확인 받은 후에만 진행. "
            "셀레니움 test_obj_rename_cascade.py 로 8 phase 자동 검증.")


# ━━━ 16. DDL 다운로드 ━━━
content_slide(16, "05-3  ·  DDL EXPORT", "DDL 다운로드",
              "방언 자동  ·  PostgreSQL  ·  Oracle  선택",
              [
                  "방언 선택 (이번 세션 v-menu 드롭다운)",
                  (1, "자동  —  모델의 데이터소스 DBMS 타입 기준"),
                  (1, "PostgreSQL  —  varchar / integer / timestamp"),
                  (1, "Oracle  —  varchar2 / number / date"),
                  (1, "물리 미연결 시 oracle 폴백"),
                  "포함 항목",
                  (1, "CREATE TABLE (모든 컬럼 + COMMENT)"),
                  (1, "PK / FK / UK 제약조건  ·  INDEX 정의"),
                  "출력",
                  (1, ".sql 파일 즉시 다운로드"),
              ], accent=CYAN)
add_notes(prs.slides[15], "DDL 다운로드는 모델 행에서 [DDL 다운로드] 드롭다운 클릭하면 방언 선택. 자동 선택하면 모델의 DBMS 타입에 맞춰 "
            "Oracle 이면 varchar2/number, Postgres 면 varchar/integer 로 출력. CREATE TABLE 만 아니라 PK/FK/UK + 인덱스 + COMMENT 까지 통합. "
            "데이터 모델 변경 후 DBA 한테 전달할 SQL 한 파일.")


# ━━━ 17. 진단 제외 (79번) ━━━
content_slide(17, "05-4  ·  DIAG TARGET  (79번 NEW)",
              "진단 제외 관리 — 79번",
              "임시 · 폐기 테이블/컬럼 → 진단 모수에서 명시적 제외 + cascade + 사유",
              [
                  "화면 구성",
                  (1, "테이블 단위 / 컬럼 단위 탭"),
                  (1, "각 행에 표준 / 구조 / 품질 3개 진단별 ON/OFF 토글 아이콘"),
                  "OFF 시 사유 모달",
                  (1, "선택 입력  —  빈 사유로도 OFF 가능 (NULL 저장)"),
                  (1, "마지막 변경자 / 일시 자동 기록"),
                  "Cascade 정책",
                  (1, "테이블 OFF  →  그 테이블의 모든 컬럼 자동 모수 제외"),
                  (1, "테이블 ON + 컬럼 OFF  →  해당 컬럼만 제외"),
                  "일괄 토글 / 권한",
                  (1, "다중 선택 → [표준 OFF / 구조 OFF / 품질 OFF / ...ON]"),
                  (1, "관리자 전용  —  일반 사용자 토글 시 403"),
              ], accent=PINK)
add_notes(prs.slides[16], "이번 세션 신규 메뉴입니다. 임시 테이블, 폐기 예정 테이블, 의미 없는 컬럼을 진단 모수에서 명시적으로 제외하는 기능. "
            "기존엔 진단 결과를 그냥 무시했는데 그러면 표준 준수율이 왜곡됩니다. 명시적 OFF 하면 모수에서 빠지므로 정확한 준수율. "
            "Cascade 가 핵심 — 테이블 OFF 하면 그 안의 컬럼 다 자동 빠지고, 테이블 ON + 컬럼만 OFF 하면 그 컬럼만. "
            "일괄 토글로 한 번에 여러 행 OFF 가능. 모든 변경은 사유와 함께 마지막 변경자/일시 기록.")


# ━━━ 18. 79번 12 phase ━━━
two_col_slide(18, "05-4-2  ·  79번 VERIFICATION", "79번 12 phase 자동 검증",
              "test_diag_target_imsi.py — Oracle DDL + 매퍼 SQL 직접 검증 / 229초 PASS",
              left_items=[
                  "Setup",
                  (1, "P1: Oracle IMSI_TEST_001/002/003 생성"),
                  (1, "P2: 메타 INSERT (12 ATTR)"),
                  "OFF 토글",
                  (1, "P3: OBJ 단건 OFF — TEST_001 표준 + 사유"),
                  (1, "P4: OBJ 일괄 OFF — TEST_002/003 구조"),
                  (1, "P5: ATTR 단건 OFF — TEST_001.NAME (사유 빈칸)"),
                  (1, "P6: ATTR 일괄 OFF — TEST_002.CODE/VALUE"),
              ],
              right_items=[
                  "매퍼 검증",
                  (1, "P7: 표준 모수 5 / 제외 7 (cascade)"),
                  (1, "P8: 구조 모수 4"),
                  "ALTER 후 재검증",
                  (1, "P9: ALTER (XYZ_DATA 길이 / NEW_COL / ETC)"),
                  (1, "P10: OFF 변경은 모수 미포함"),
                  "복귀 + Cleanup",
                  (1, "P11: ON 복귀 후 모수 13 / 사유 NULL"),
                  (1, "P12: cleanup — Oracle DROP + 메타 DELETE"),
              ],
              left_label="P1 — P6 / SETUP + OFF",
              right_label="P7 — P12 / VERIFY + CLEANUP",
              accent=PINK, accent2=VIOLET)
add_notes(prs.slides[17], "79번 검증은 12 phase 로 끊어서 매퍼 SQL 까지 직접 검증합니다. "
            "Oracle 컨테이너에 임시 테이블 만들고 메타 INSERT 한 다음 단건/일괄 OFF 적용. "
            "그 다음 표준 진단 매퍼와 구조 진단 매퍼가 OFF row 를 정확히 모수에서 빼는지 SQL 수준에서 카운트 확인. "
            "ALTER 후 재검증 단계는 OFF 처리된 테이블의 변경이 결과에 안 떠야 한다는 핵심 동작 검증. "
            "마지막에 ON 복귀 후 모수 회복까지 풀 사이클 — 229초 PASS.")


# ━━━ 19. 표준화 진단 실행 ━━━
content_slide(19, "06-1  ·  STND DIAG  ·  RUN", "표준화 진단 — 실행",
              "q-executor 백그라운드  +  STOMP 실시간 진행률",
              [
                  "실행 절차",
                  (1, "[표준 진단 > 진단 실행] → 모델 선택 → [진단 시작]"),
                  (1, "q-executor 가 백그라운드 처리 (q-center 와 STOMP 통신)"),
                  "이슈 6종",
                  (1, "용어 미존재 / 한글명 불일치 / 영문약어 불일치"),
                  (1, "타입 불일치 / 길이 불일치 / 도메인 불일치"),
                  "진행률 표시",
                  (1, "현재 처리 중 / 전체 컬럼수 / 발견된 이슈 / 예상 완료 시각"),
                  "결과 저장",
                  (1, "TB_DIAG_JOB (헤더) + TB_DIAG_RESULT (이슈 건별)"),
                  (1, "TARGET_YN='Y' 만 모수  ·  79번 진단 제외 cascade 즉시 반영"),
              ], accent=CYAN)
add_notes(prs.slides[18], "표준 진단은 사용자가 모델 선택하고 [진단 시작] 누르면 q-executor 가 백그라운드에서 처리. "
            "진행률은 WebSocket STOMP 로 실시간 push 되니까 사용자는 화면 닫아도 됩니다. "
            "이슈 6종 — 용어 미존재, 한글명 불일치, 영문약어 불일치, 타입 불일치, 길이 불일치, 도메인 불일치. "
            "결과는 TB_DIAG_JOB 헤더 + TB_DIAG_RESULT 이슈별 row 로 저장. 79번 진단 제외 cascade 가 매퍼 단계에서 즉시 반영됩니다.")


# ━━━ 20. 표준화 진단 결과 ━━━
content_slide(20, "06-2  ·  STND DIAG  ·  RESULT", "표준화 진단 — 결과",
              "필터  +  상세 drawer  +  표준 준수율",
              [
                  "필터",
                  (1, "이슈 유형 / 테이블 / 컬럼 / 표준 준수 여부 / 확인 여부"),
                  "그리드",
                  (1, "행 클릭 → 상세 drawer (현재값 vs 권장값 비교)"),
                  "표준 준수율",
                  (1, "(전체 컬럼수 − 이슈 컬럼수) / 전체 컬럼수 × 100"),
                  (1, "RESULT_CNT(이슈 건수) 가 아닌 ISSUE_COL_CNT(이슈 컬럼수) 사용"),
                  (1, "한 컬럼이 6 이슈라도 1로 카운트  —  의미 있는 백분율"),
                  "결과 다운로드",
                  (1, "현재 필터 상태 그대로 POI 동적 XLSX"),
              ], accent=CYAN)
add_notes(prs.slides[19], "진단 결과 화면은 필터로 빠르게 좁혀들어갑니다. 행 클릭하면 우측에 drawer 가 열려서 현재값과 권장값을 나란히 비교. "
            "표준 준수율 계산식이 중요한데, 이슈 건수가 아니라 이슈 컬럼 수로 계산합니다. "
            "한 컬럼이 6개 이슈가 있어도 1로 카운트 — 그래야 의미있는 백분율이 나오고 사용자가 컬럼 단위로 작업 우선순위를 잡을 수 있어요.")


# ━━━ 21. 빠른 등록 + 코멘트 ━━━
content_slide(21, "06-3  ·  QUICK REGISTER", "진단 결과 — 용어 빠른 등록 + 코멘트",
              "이슈 → 즉시 표준 사전 보강 → 진단 재실행 없이 결과 갱신",
              [
                  "용어 빠른 등록",
                  (1, "이슈 행 → [용어 빠른 등록] → 즉시 등록"),
                  (1, "진단 재실행 없이 결과 row 갱신 (TB_DIAG_RESULT 한 row update)"),
                  "코멘트",
                  (1, "진단 결과에 메모 (예: '검토 완료', '현재 시스템 제약상 불가')"),
                  (1, "변경 이력에 자동 등록  —  누가 언제 코멘트했는지 추적"),
                  "[해결] 버튼",
                  (1, "PARTIAL/FAILED 행 → 자동 표준화 추천 모달 진입"),
                  (1, "53번 컬럼 그리드와 동일 모달 재사용  —  학습 비용 0"),
              ], accent=CYAN)
add_notes(prs.slides[20], "진단 결과에서 이슈 행마다 [용어 빠른 등록] 버튼이 있어요. 표준 사전에 없는 용어가 발견되면 그 자리에서 등록해버리고 "
            "진단 재실행 없이 결과 row 만 갱신. 운영 중 이걸로 빠르게 표준 사전 채워나갑니다. "
            "코멘트는 진단 결과에 메모를 남길 수 있고, 누가 언제 코멘트했는지 변경 이력에 자동 등록. "
            "[해결] 버튼은 자동 표준화 추천 모달로 점프 — 53번 컬럼 그리드의 모달과 동일한 모달이라 학습 비용 없습니다.")


# ━━━ 22. 구조 변경 진단 ━━━
content_slide(22, "07  ·  STRUCT DIAG", "구조 변경 진단",
              "DBMS 실 스키마 vs 수집 스냅샷  —  추가/변경/삭제 자동 감지",
              [
                  "실행",
                  (1, "[구조 변경 진단 > 진단 실행]  →  모델 선택  →  [실행]"),
                  (1, "물리 모델만 (DSID 있는)  —  논리 모델은 거부 (400)"),
                  "결과 화면 3단계 조회",
                  (1, "이력 (TB_STRUCT_DIAG_HISTORY)  →  상세 (DETAIL)  →  컬럼/제약/인덱스별"),
                  (1, "각 변경 — 컬럼명, 변경 전/후 타입·길이·NULL 여부"),
                  "79번 진단 제외 cascade",
                  (1, "OFF 표시된 테이블/컬럼의 변경은 결과 미등장"),
                  (1, "PC1 보강 (2bef4e8)  —  prev/curr OFF set 통합 + toUpperCase 정렬"),
              ], accent=CYAN)
add_notes(prs.slides[21], "구조 변경 진단은 데이터베이스의 실제 스키마와 우리가 수집한 스냅샷을 비교해서 컬럼이 추가/변경/삭제됐는지 감지. "
            "결과는 3단계로 조회 — 진단 이력, 그 안의 상세, 그 안의 컬럼/제약/인덱스별. "
            "79번 진단 제외 cascade 가 여기서도 적용 — OFF 처리된 테이블이나 컬럼의 변경은 결과에 안 떠서 노이즈가 줄어듭니다. "
            "이 cascade 통합 작업은 어제 PC1 에서 보강한 부분이고 4 케이스 전수 검증됐습니다.")


# ━━━ 23. 자동 표준화 분석 ━━━
content_slide(23, "08-1  ·  AUTO STANDARDIZATION", "자동 표준화 — 컬럼 표준화 분석",
              "한글 컬럼명  →  단어 분리  →  영문약어 / 도메인 추천  →  자동 등록",
              [
                  "입력",
                  (1, "[자동 표준화 지원 > 컬럼 표준화]  →  textarea 줄바꿈 구분"),
                  (1, "공백 무시 / 중복 제거"),
                  "분석 알고리즘",
                  (1, "DP 점수  —  SPLIT_PENALTY=8000 / 1자 비형식단어 격하(5000)"),
                  (1, "사후처리 resolveUncertainRuns  —  미신뢰 토큰 재검색"),
                  (1, "동의어 매핑 synmToWord  —  카테고리 → 범주 자동 매칭"),
                  "결과",
                  (1, "REGISTERED / AUTO / PARTIAL / FAILED  4 status 자동 판정"),
                  (1, "추천 영문약어 / 도메인 / 데이터타입 / 길이 자동 채움"),
              ], accent=PINK)
add_notes(prs.slides[22], "자동 표준화는 DataQ 의 핵심 차별화 영역입니다. 한글 컬럼명을 textarea 에 줄바꿈 구분해서 여러 건 동시 입력 → 분석. "
            "알고리즘은 DP 기반으로 가장 그럴듯한 단어 분리를 찾고, DP 만으로 부족한 미신뢰 토큰은 사후처리로 사전 재검색해서 보강. "
            "동의어 매핑이 결정적인데, 예를 들어 '카테고리'를 입력하면 알고리즘이 '범주' 단어로 자동 매칭해줍니다. "
            "결과는 4가지 status 로 자동 분류되고 추천 영문약어/도메인/타입/길이까지 채워줍니다.")


# ━━━ 24. 수정 모달 ━━━
content_slide(24, "08-2  ·  EDIT MODAL", "자동 표준화 — 수정 모달 (cascade)",
              "행안부 지침 — '분류어' → '형식단어' 가시 텍스트 통일",
              [
                  "모달 구조",
                  (1, "단어 테이블  —  자동 분리 결과 + 수동 추가/제거"),
                  (1, "형식단어 검색·자동완성 + [형식단어 추가] 버튼"),
                  (1, "용어 도메인  —  마지막 단어의 분류명 cascade"),
                  (1, "용어 미리보기  —  입력값으로 만들어질 한글/영문 즉시 표시"),
                  "Cascade 동작",
                  (1, "단어 테이블 변경 → 마지막 단어의 분류명 → 도메인 후보 자동 재로드"),
                  (1, "마지막 단어가 형식단어가 아니면 도메인 비활성 + 빨간 안내"),
                  "검증",
                  (1, "test_ca8858d_clsf_domain.py — 10/10 PASS"),
              ], accent=PINK)
add_notes(prs.slides[23], "이 수정 모달이 우리 시스템의 UX 차별화 중 하나예요. PARTIAL 이나 FAILED 행에서 [수정] 누르면 진입. "
            "단어 테이블에서 형식단어를 검색·자동완성으로 골라서 추가하면 마지막 단어의 분류명에 따라 도메인 드롭다운이 자동으로 cascade 됩니다. "
            "마지막 단어가 형식단어 (분류어) 가 아니면 도메인이 비활성 + 빨간 안내. "
            "한글/영문 미리보기가 실시간으로 갱신되니까 사용자가 결과를 즉시 확인. test_ca8858d_clsf_domain 10 케이스 PASS.")


# ━━━ 25. 데이터 품질 검증항목 ━━━
content_slide(25, "09-1  ·  QUAL  ·  RULES", "데이터 품질 진단 — 검증항목 (룰 정의)",
              "DQI / CTQ / BR  +  도메인 룰 1:N  +  컬럼 매핑",
              [
                  "룰 종류",
                  (1, "품질지표 (DQI)  —  데이터 품질 지표"),
                  (1, "핵심관리항목 (CTQ)  —  비즈니스 우선순위 컬럼"),
                  (1, "업무규칙 (BR)  —  도메인 정합성"),
                  "룰 카탈로그 (TB_QUAL_RULE_CATALOG)",
                  (1, "이메일 / 주민번호 / 사업자번호 / 휴대전화번호 등 표준 정규식"),
                  "도메인 룰 (1:N)",
                  (1, "한 도메인에 여러 룰 매핑 (예: 이메일 + 정규식 + NOT NULL)"),
                  "컬럼별 룰 매핑",
                  (1, "effective rule = 도메인 룰 ∪ 컬럼 직접 매핑  —  1 SQL JOIN"),
              ], accent=PINK)
add_notes(prs.slides[24], "데이터 품질 진단은 67/70번 설계의 결과로 신규 추가된 영역입니다. 룰을 도메인 단위로 정의해서 같은 도메인 컬럼에 일괄 적용하고, "
            "컬럼별로 룰을 직접 매핑할 수도 있습니다. effective rule 은 도메인 룰과 컬럼 매핑을 합집합으로 1 SQL 에 JOIN. "
            "룰 카탈로그에는 이메일·주민번호 같은 표준 정규식이 미리 들어있어서 즉시 사용 가능.")


# ━━━ 26. 검증대상 + 품질검증 ━━━
content_slide(26, "09-2  ·  QUAL  ·  RUN", "검증대상 + 품질검증 실행",
              "컬럼 매핑  +  값 프로파일  +  룰 위반 동시",
              [
                  "검증대상 (DSQualColRule)",
                  (1, "컬럼 그리드 → 효과 룰 (effective rule) 표시"),
                  (1, "단위 재진단  —  한 컬럼만 즉시 재실행"),
                  "품질검증 실행 (DSQualValueProfile)",
                  (1, "모델 선택  →  [DB 연결된 모델만] 자동 (connectedOnly='Y')"),
                  (1, "컬럼 체크  →  [진단 시작]"),
                  "처리 (q-executor)",
                  (1, "값 프로파일  —  null/distinct/min/max/avg/quantile"),
                  (1, "룰 위반 검사  —  RuleSqlBuilder dialect (Postgres/Oracle/Cubrid)"),
                  "결과 저장",
                  (1, "TB_QUAL_PROFILE_RESULT (UPSERT) + TB_QUAL_RULE_RESULT + TB_QUAL_VIOLATION_SAMPLE"),
              ], accent=PINK)
add_notes(prs.slides[25], "검증대상 화면에서 컬럼별로 어떤 룰이 적용되는지 효과 룰을 확인하고, 단위 재진단으로 한 컬럼만 즉시 재실행 가능. "
            "품질검증 실행은 DB 연결된 모델만 자동 필터되니까 논리 모델은 안 보입니다. "
            "컬럼 체크해서 진단 시작하면 q-executor 가 값 프로파일과 룰 위반 검사를 동시 수행. "
            "RuleSqlBuilder 가 dialect 분기 — Postgres/Oracle/Cubrid 별로 정규식 표현이 달라서 그 부분 호환성 직접 작업.")


# ━━━ 27. 결과 + 시계열 ━━━
content_slide(27, "09-3  ·  QUAL  ·  RESULT", "데이터 품질 진단 — 결과 + 시계열",
              "위반율 / 위반 샘플 / 검증항목별 / 시계열 추이",
              [
                  "품질검증 결과",
                  (1, "진단별 위반율 + 위반 샘플 (PK + 위반값, 룰당 100건 default)"),
                  (1, "drawer 상세  —  룰 텍스트 + 위반 행 그리드"),
                  "테이블별 결과 / 검증항목별 결과",
                  (1, "테이블 단위 합계 + 컬럼별 위반율"),
                  (1, "룰 단위  —  어떤 룰이 어디서 위반됐는지"),
                  "데이터 품질 현황",
                  (1, "시계열  —  진단 N회의 위반율 추이 차트"),
                  "79번 진단 제외 (컬럼 단위)",
                  (1, "QUAL_DIAG_TARGET_YN='N' 컬럼은 모수 자동 제외"),
                  (1, "테이블 단위 cascade 는 67/70번 정식 통합 시 보강 예정"),
              ], accent=PINK)
add_notes(prs.slides[26], "품질검증 결과 화면은 4가지 관점 — 진단 단위, 테이블 단위, 검증항목(룰) 단위, 시계열 추이. "
            "위반 샘플은 PK 와 위반값을 같이 보여주니까 운영자가 어떤 row 가 문제인지 즉시 확인 가능. 룰당 기본 100건. "
            "79번 진단 제외는 컬럼 단위로 즉시 적용되고, 테이블 단위 cascade 는 67/70번 정식 통합 시 보강 예정으로 메모돼있습니다.")


# ━━━ 28. 진단 스케줄 등록 ━━━
content_slide(28, "10-1  ·  SCHEDULE", "진단 스케줄 — 등록",
              "cron 표현식  +  진단 유형  +  즉시 실행  +  동시 실행 방어",
              [
                  "스케줄 등록",
                  (1, "[진단 스케줄] → [등록] 모달"),
                  (1, "cron 표현식 (예: '0 0 9 * * MON-FRI'  —  평일 9시)"),
                  (1, "진단 유형 — STND / STRUCT / BOTH"),
                  (1, "대상 모델 + 활성/비활성 토글"),
                  "즉시 실행 (runNow)",
                  (1, "테스트용 강제 트리거  —  cron 무시 + 즉시 결과 row 생성"),
                  "동시 실행 방어",
                  (1, "같은 모델/유형 진행 중이면 SKIP (자동 검증)"),
                  "권한",
                  (1, "등록·수정·삭제·즉시 실행은 관리자 전용 (admin gate)"),
              ], accent=VIOLET)
add_notes(prs.slides[27], "진단 스케줄러는 Phase 1~4 로 단계 구현했고 자동화 운영을 위한 핵심입니다. cron 표현식으로 스케줄 등록, "
            "진단 유형은 표준만/구조만/둘 다 셋 중 선택. 대상 모델 지정. 활성/비활성 토글로 일시 중단 가능. "
            "테스트할 때는 cron 기다리지 말고 즉시 실행 버튼으로 강제 트리거. "
            "같은 모델 같은 유형이 이미 돌고 있으면 자동으로 SKIP — 동시 실행 방어. 셀레니움으로 SKIP 정책까지 자동 검증.")


# ━━━ 29. 스케줄 실행 이력 ━━━
content_slide(29, "10-2  ·  SCHEDULE LOG", "스케줄 실행 이력",
              "시간 / 모델 / 유형 / 상태 / 소요시간  +  drawer 상세",
              [
                  "그리드 컬럼",
                  (1, "실행 시각 / 모델명 / 진단 유형 / 상태 (성공/실패/SKIP)"),
                  (1, "소요 시간 / 실행 사용자 / 결과 카운트"),
                  "필터",
                  (1, "기간 (from/to) / 모델 / 유형 / 상태"),
                  "Drawer 상세",
                  (1, "행 클릭 → 상세 — 진단 ID 링크 (결과 화면으로 점프)"),
                  (1, "에러 케이스  —  스택트레이스 + 원인 분석 단서"),
              ], accent=VIOLET)
add_notes(prs.slides[28], "스케줄 실행 이력은 자동 진단의 운영 로그. 시각/모델/유형/상태/소요시간 전부 한 그리드. "
            "필터로 기간이나 상태별로 좁혀 보고, 행 클릭하면 drawer 가 열려서 진단 ID 링크로 진단 결과 화면 바로 점프. "
            "실패 케이스에는 에러 스택트레이스가 들어있어서 원인 추적 가능합니다.")


# ━━━ 30. 권한 ━━━
content_slide(30, "10-3  ·  PERM GATE", "진단 스케줄 — 권한 (관리자 게이트)",
              "일반 사용자는 조회만  —  UI 비활성 + API 403 양쪽 검증",
              [
                  "관리자 (isAdmin=true)",
                  (1, "[등록] / [수정] / [삭제] / [즉시 실행] 버튼 노출"),
                  (1, "스케줄 활성/비활성 토글 가능"),
                  "일반 사용자",
                  (1, "그리드 조회만  —  모든 액션 버튼 비활성"),
                  (1, "API 직접 호출 시도  →  403 Forbidden (서버 검증)"),
                  "검증",
                  (1, "test_phase4_ui_admin_gate.py  —  4/4 PASS"),
                  (1, "test_perm_matrix.py  —  admin/user 권한 매트릭스 자동 검증"),
              ], accent=VIOLET)
add_notes(prs.slides[29], "스케줄 자체는 일반 사용자도 조회 가능 (운영 투명성). 다만 등록/수정/삭제/즉시 실행은 관리자 전용. "
            "UI 비활성만으로는 부족하니까 API 도 서버 단에서 isAdmin 체크 후 403. "
            "셀레니움 phase4_ui_admin_gate 와 perm_matrix 두 테스트로 권한 매트릭스 자동 검증 — 회귀 즉시 감지.")


# ━━━ 31. 마이페이지 ━━━
content_slide(31, "11-1  ·  MY PAGE", "마이페이지",
              "내 정보 + 요청 현황",
              [
                  "내 정보",
                  (1, "비밀번호 변경 (현재 비밀번호 검증 + 새 비밀번호 + 확인)"),
                  (1, "마지막 로그인 시각"),
                  "요청 현황 (DSMyRequest)",
                  (1, "본인이 신청한 단어/용어/도메인 승인 상태"),
                  (1, "카드 4개  —  전체 / 대기 / 승인 / 반려 (필터 클릭 가능)"),
                  (1, "그리드 + 상세 패널 (반려 사유 표시)"),
                  (1, "검색  —  유형 / 기간 (datetime range)"),
              ], accent=INDIGO)
add_notes(prs.slides[30], "마이페이지는 사용자 본인이 신청한 항목들을 한눈에 보는 화면. 카드 4개에 전체/대기/승인/반려 카운트가 보이고 클릭하면 필터. "
            "행 클릭하면 상세 패널이 열려서 반려된 경우 사유까지 표시. 운영자가 자기 요청 진척도 빠르게 확인할 수 있어요.")


# ━━━ 32. 관리 ━━━
content_slide(32, "11-2  ·  ADMIN", "관리 (관리자 전용)",
              "사용자 + 승인 + 데이터 소스",
              [
                  "사용자 관리",
                  (1, "사용자 등록 / 비활성화 / 권한 변경 (admin 토글)"),
                  (1, "비밀번호 초기화 (관리자 강제 리셋)"),
                  "승인",
                  (1, "행별 [승인]/[반려] 인라인 (사유 입력)"),
                  (1, "단어 선승인 + cascade 반려 + 반려 후 물리 삭제 정책"),
                  "데이터 소스",
                  (1, "외부 DBMS 연결 정보 (jasypt 암호화)"),
                  (1, "Oracle SID/Service Name 양쪽 지원 (drivers.xml swap)"),
                  (1, "[연결 테스트]  —  즉시 검증"),
                  "메뉴 가시성",
                  (1, "isAdmin=false  →  '관리' 그룹 자체 DOM 미렌더 (서버 + 클라이언트)"),
              ], accent=INDIGO)
add_notes(prs.slides[31], "관리 메뉴는 관리자 전용. 사용자 관리 / 승인 / 데이터 소스 3개. 데이터 소스가 핵심인데 외부 DBMS 연결 정보 등록하고 "
            "[연결 테스트] 로 즉시 확인. 비밀번호는 jasypt 로 암호화 저장. Oracle 은 SID 와 Service Name 둘 다 지원. "
            "관리 그룹은 isAdmin 체크해서 일반 사용자 화면에는 DOM 자체가 안 만들어집니다 — 클라이언트와 서버 양쪽 검증.")


# ━━━ 33. 마스터 러너 ━━━
content_slide(33, "12-1  ·  AUTOMATION", "셀레니움 자동화 회귀 — 마스터 러너",
              "29건 통합 테스트 자동 회귀  ·  ~28분  ·  종료코드 0",
              [
                  "1 명령어 전체 실행",
                  (1, "python  dataQ설계/테스트/selenium/run_all.py"),
                  "그룹별 진행",
                  (1, "API/Login 5  ·  표준 사전 5  ·  데이터 모델 11"),
                  (1, "논리물리 3  ·  진단 스케줄 4  ·  79번 1"),
                  "Per-test cleanup hook",
                  (1, "DB 폴루션 자동 DELETE  —  ^(셀|테스트), 셀도메인%, IMSI_*"),
                  (1, "좀비 Edge 프로세스 정리  —  taskkill msedgedriver.exe"),
                  (1, "→ 28건 순차 실행 시 누적 폴루션으로 인한 플레이키 0"),
                  "변경 후 즉시 회귀",
                  (1, "각 테스트 단일 실행 가능  —  변경 영향 영역만 빠르게 검증"),
              ], accent=SUCCESS)
add_notes(prs.slides[32], "셀레니움 통합 테스트 29건이 한 명령어로 자동 회귀. 약 28분 걸립니다. 종료코드 0 이면 모두 PASS. "
            "각 테스트 사이에 cleanup hook 이 DB 폴루션 자동 삭제하고 좀비 Edge 프로세스 정리 — 이걸 안 하면 누적 데이터로 플레이키 발생. "
            "변경 후엔 영향 영역의 단일 테스트만 빠르게 돌리고, 큰 변경 후엔 전체 회귀. 회귀 즉시 감지 패턴.")


# ━━━ 34. 회귀 결과 (KPI 카드) ━━━
stat_slide(34, "12-2  ·  KPI", "회귀 결과 — 29 PASS / 0 FAIL",
           "동일 환경 2회 연속 동일 결과  ·  cleanup hook 안정성 확인",
           stats=[
               {'value': '29', 'label': 'TOTAL TESTS', 'color': INDIGO},
               {'value': '29', 'label': 'PASS', 'color': SUCCESS},
               {'value': '0',  'label': 'FAIL', 'color': DANGER},
               {'value': '1675s', 'label': 'ELAPSED', 'color': VIOLET},
               {'value': '±50s', 'label': 'STDDEV (2회)', 'color': CYAN},
           ],
           items=[
               "그룹별 결과",
               (1, "API/Login (가벼움)  —  5/5 PASS"),
               (1, "표준 사전  —  5/5 PASS (test_term_register_v2 v2 포함)"),
               (1, "데이터 모델 (Phase 5)  —  11/11 PASS"),
               (1, "논리/물리 모델 진단  —  3/3 PASS"),
               (1, "진단 스케줄 (Phase 2~4)  —  4/4 PASS"),
               (1, "진단 제외 관리 (79번)  —  1/1 PASS (12 phase, 228초)"),
           ],
           accent=SUCCESS, accent2=CYAN)
add_notes(prs.slides[33], "오늘 마스터 러너 회귀 — 29건 PASS 0 FAIL. 1675초 약 28분. 79번 진단 제외 관리는 12 phase 가 있어서 단일 테스트가 228초. "
            "표준 사전 그룹에 81/82번 용어 등록 v2 가 5번째로 추가됐고 8 케이스 PASS. "
            "동일 환경에서 2회 연속 동일 결과 — 1624초 / 1675초 (편차 ±50초). cleanup hook 으로 재현성 확보.")


# ━━━ 35. SMETA 비교 ━━━
two_col_slide(35, "13  ·  SMETA COMPARE", "SMETA 대비 차별화",
              "75/76/77/78번 분석 결과  —  자동 표준화 + 진단 제외 + 권한 자동 검증 + 한국 표준",
              left_items=[
                  "자동 표준화 엔진",
                  (1, "DataQ — DP 점수 + 동의어 cascade + 사후처리 재검색 (자동)"),
                  (1, "SMETA — 단어 분리 수동, 영문약어 사용자 직접 입력"),
                  "진단 제외 관리",
                  (1, "DataQ — 메뉴 + cascade + 사유 + 변경 이력 (이번 세션 신규)"),
                  (1, "SMETA — 미지원"),
              ],
              right_items=[
                  "권한 매트릭스 자동 검증",
                  (1, "DataQ — test_perm_matrix / phase4_ui_admin_gate 자동 회귀"),
                  (1, "SMETA — 수동 검증"),
                  "한국 표준 직접 반영",
                  (1, "DataQ — 행안부 형식단어, 분류어→형식단어, 표준 단어 3,300건"),
                  (1, "SMETA — 일반 표준 (한국 특화 보강 별도)"),
              ],
              left_label="DataQ STRENGTHS  ·  1, 2",
              right_label="DataQ STRENGTHS  ·  3, 4",
              accent=PINK, accent2=VIOLET)
add_notes(prs.slides[34], "75/76/77/78번 4개 문서로 SMETA 와 정밀 비교한 결과 4가지 차별화 영역. "
            "자동 표준화 엔진이 가장 큰 차이 — DataQ 는 한글 입력만으로 단어 분리, 영문약어, 도메인까지 자동인데 SMETA 는 수동. "
            "진단 제외 관리는 DataQ 만의 신규 기능. 권한 매트릭스도 DataQ 는 자동 셀레니움 회귀. "
            "그리고 한국 표준 — 행안부 형식단어 개념과 표준 단어 3,300건 시드까지 직접 반영했습니다.")


# ━━━ 36. 마무리 ━━━
s = add_slide()
add_bg(s, DARK_BG)
add_accent_blob(s, Inches(8.5), Inches(-2.0), Inches(8.0).emu, VIOLET, alpha=22000)
add_accent_blob(s, Inches(-3.0), Inches(3.5), Inches(7.0).emu, CYAN, alpha=18000)
add_accent_blob(s, Inches(9.0), Inches(5.0), Inches(5.0).emu, PINK, alpha=12000)
# 라벨
add_text(s, Inches(0.7), Inches(0.6), Inches(6.0), Inches(0.4),
         "CLOSING", size=12, bold=True, color=CYAN, font='Consolas')
# 타이틀
add_text(s, Inches(0.7), Inches(1.8), Inches(12.0), Inches(1.5),
         "Thank you", size=88, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(3.5), Inches(12.0), Inches(0.7),
         "Q & A", size=36, color=RGBColor(0xC7, 0xD2, 0xFE))
# 액센트
add_rect(s, Inches(0.7), Inches(4.45), Inches(0.8), Pt(5), VIOLET)
add_rect(s, Inches(1.6), Inches(4.45), Inches(0.5), Pt(5), CYAN)
# 요약 카드
card = slide_card = add_card(s, Inches(0.7), Inches(4.85), Inches(12.0), Inches(1.95),
                              fill=DARK_BG2, border=DARK_BG2)
add_text(s, Inches(1.0), Inches(5.0), Inches(11.5), Inches(0.4),
         "오늘 시연 요약", size=12, bold=True, color=CYAN, font='Consolas')
add_text(s, Inches(1.0), Inches(5.4), Inches(11.5), Inches(1.4),
         ["·  좌측 메뉴 1~12 영역 핵심 기능 + 신규 (79번 진단 제외, 81/82번 용어 등록 v2)",
          "·  자동화 회귀  29 PASS / 0 FAIL  /  1675초",
          "·  SMETA 대비 차별화  4 영역  (자동 표준화, 진단 제외, 권한 자동 검증, 한국 표준)"],
         size=14, color=WHITE, line_spacing=1.5)
add_text(s, Inches(0.4), Inches(7.0), Inches(8.0), Inches(0.4),
         "Narae DataQ · 종합 시연 2026-05-06 · 발표 — 장재영", size=10, color=TEXT_DIM)
add_text(s, Inches(12.4), Inches(7.0), Inches(0.8), Inches(0.4),
         "36 / 36", size=10, color=TEXT_DIM, align=PP_ALIGN.RIGHT, font='Consolas')
add_notes(s, "이상 Narae DataQ 종합 시연 마칩니다. 오늘 좌측 메뉴 순서 그대로 모든 핵심 기능과 이번 세션 신규 기능 "
            "(79번 진단 제외, 81/82번 용어 등록 v2) 보여드렸고, 셀레니움 자동 회귀로 29 PASS / 0 FAIL 까지 검증된 상태입니다. "
            "SMETA 대비 차별화 4 영역도 정리. 질문 받겠습니다.")


# 저장
prs.save(OUT)
print(f"PPT 생성: {OUT}")
print(f"슬라이드 수: {len(prs.slides)}")
